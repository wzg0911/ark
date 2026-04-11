#!/usr/bin/env /usr/bin/python3
# -*- coding: utf-8 -*-
"""
元神项目融资BP生成脚本
基于《元神项目核心知识库》生成10页专业PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ============================================================
# 颜色主题（深色风格 + Notion风）
# ============================================================
COLOR_BG_DARK     = RGBColor(0x0D, 0x0D, 0x1A)
COLOR_BG_CARD     = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT      = RGBColor(0x6C, 0x63, 0xFF)
COLOR_ACCENT2     = RGBColor(0x00, 0xD4, 0xAA)
COLOR_TEXT_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT_GRAY   = RGBColor(0xA0, 0xA0, 0xB0)
COLOR_PLACEHOLDER = RGBColor(0x2A, 0x2A, 0x45)
COLOR_DIVIDER     = RGBColor(0x3A, 0x3A, 0x5C)
COLOR_GOLD        = RGBColor(0xF0, 0xC0, 0x40)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ============================================================
# 工具函数
# ============================================================

def set_bg(slide, prs, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=COLOR_TEXT_WHITE,
                 align=PP_ALIGN.LEFT, font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_phbox(slide, left, top, width, height, label="图片占位"):
    add_rect(slide, left, top, width, height, COLOR_PLACEHOLDER)
    txBox = slide.shapes.add_textbox(left, top + height/2 - Inches(0.15),
                                      width, Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {label} ]"
    run.font.size = Pt(10)
    run.font.color.rgb = COLOR_TEXT_GRAY


def add_divider(slide, left, top, width, color=COLOR_ACCENT):
    line = slide.shapes.add_shape(1, left, top, width, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_bullet_card(slide, title, items, left, top, width, height, accent):
    add_rect(slide, left, top, width, height, COLOR_BG_CARD)
    add_rect(slide, left, top, Inches(0.06), height, accent)
    add_text_box(slide, title, left+Inches(0.15), top+Inches(0.08),
                 width-Inches(0.2), Inches(0.35), font_size=13, bold=True, color=accent)
    y = top + Inches(0.42)
    for item in items:
        add_text_box(slide, f"• {item}", left+Inches(0.15), y,
                     width-Inches(0.2), Inches(0.25), font_size=11, color=COLOR_TEXT_WHITE)
        y += Inches(0.26)


def add_page_num(slide, n, total=10):
    add_text_box(slide, f"{n} / {total}", SLIDE_W-Inches(1.2),
                 SLIDE_H-Inches(0.42), Inches(1), Inches(0.28),
                 font_size=9, color=COLOR_TEXT_GRAY, align=PP_ALIGN.RIGHT)


def add_footer(slide, text="元神项目 · 融资计划书 · 2026"):
    add_text_box(slide, text, Inches(0.4), SLIDE_H-Inches(0.42),
                 Inches(5), Inches(0.28), font_size=8, color=COLOR_TEXT_GRAY)


def add_section_title(slide, num, title_en, title_cn):
    add_text_box(slide, num, Inches(0.5), Inches(0.28), Inches(1), Inches(0.65),
                 font_size=38, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, title_en, Inches(0.5), Inches(0.82), Inches(7), Inches(0.45),
                 font_size=20, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.5), Inches(1.32), SLIDE_W-Inches(1), COLOR_DIVIDER)


# ============================================================
# 10页生成函数
# ============================================================

def page1_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    # 顶部紫线
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), COLOR_ACCENT)
    # 中心内容区
    add_divider(slide, Inches(0), Inches(2.5), SLIDE_W, COLOR_ACCENT)
    add_text_box(slide, "META SOUL",
                 Inches(0.8), Inches(2.7), Inches(11.5), Inches(1.1),
                 font_size=64, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, "元神项目",
                 Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.8),
                 font_size=36, color=COLOR_TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(5), Inches(4.55), Inches(3.3), COLOR_ACCENT2)
    add_text_box(slide, "让每个人都有一个比自己更好的自己",
                 Inches(1), Inches(4.8), Inches(11), Inches(0.55),
                 font_size=17, color=COLOR_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "个人精神世界模型构建器",
                 Inches(1), Inches(5.4), Inches(11), Inches(0.45),
                 font_size=13, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)
    # 底部
    add_rect(slide, Inches(0), SLIDE_H-Inches(0.06), SLIDE_W, Inches(0.06), COLOR_ACCENT)
    add_text_box(slide, "融资计划书 2026",
                 Inches(0.8), SLIDE_H-Inches(0.7), Inches(11.5), Inches(0.38),
                 font_size=11, color=COLOR_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_page_num(slide, 1)


def page2_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_section_title(slide, "01", "问题", "Problem")
    problems = [
        ("数字足迹碎片化", ["个人数据分散在微信/微博/小红书等平台",
                          "无法形成完整的自我认知图谱",
                          "平台割裂导致记忆断层"]),
        ("AI不理解'你'",  ["通用AI知道世界，但不知道你的故事",
                          "无法理解你的偏好结构与价值观",
                          "每个新对话都是陌生人"]),
        ("关系被平台绑架",["社交关系存储在平台，平台消失即消失",
                          "用户对个人数据无主权",
                          "数据价值被平台收割，用户零收益"]),
    ]
    cw = Inches(3.85); ch = Inches(4.3)
    colors = [COLOR_ACCENT, COLOR_ACCENT2, COLOR_GOLD]
    for i, (t, its) in enumerate(problems):
        x = Inches(0.5) + i*(cw+Inches(0.3))
        add_bullet_card(slide, t, its, x, Inches(1.55), cw, ch, colors[i])
        add_text_box(slide, f"P{i+1}", x, Inches(1.55)+ch+Inches(0.08), cw, Inches(0.22),
                     font_size=8, color=colors[i], align=PP_ALIGN.RIGHT)
    add_footer(slide); add_page_num(slide, 2)


def page3_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_section_title(slide, "02", "解决方案", "Solution")
    # 左侧卡片
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(5.9), Inches(5.3), COLOR_BG_CARD)
    add_text_box(slide, "核心价值主张", Inches(0.7), Inches(1.68), Inches(5.5), Inches(0.38),
                 font_size=13, bold=True, color=COLOR_ACCENT)
    add_divider(slide, Inches(0.7), Inches(2.06), Inches(5.5), COLOR_ACCENT)
    layers = [("第一层：分身","替你干活，让你自由",COLOR_ACCENT),
              ("第二层：高我","比你更懂你，让你更强",COLOR_ACCENT2),
              ("第三层：元神","与你合一，让你绽放",COLOR_GOLD)]
    for i,(t,d,c) in enumerate(layers):
        y = Inches(2.3)+i*Inches(1.3)
        add_rect(slide, Inches(0.7), y, Inches(0.08), Inches(0.9), c)
        add_text_box(slide, t, Inches(0.88), y, Inches(5.3), Inches(0.33),
                     font_size=13, bold=True, color=COLOR_TEXT_WHITE)
        add_text_box(slide, d, Inches(0.88), y+Inches(0.3), Inches(5.3), Inches(0.28),
                     font_size=10, color=COLOR_TEXT_GRAY)
    # 右侧卡片
    add_rect(slide, Inches(6.7), Inches(1.55), Inches(6.1), Inches(5.3), COLOR_BG_CARD)
    add_text_box(slide, "一句话定位", Inches(6.9), Inches(1.68), Inches(5.7), Inches(0.38),
                 font_size=13, bold=True, color=COLOR_ACCENT2)
    add_text_box(slide, '"让AI理解你的世界"',
                 Inches(6.9), Inches(2.15), Inches(5.7), Inches(0.55),
                 font_size=16, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(6.9), Inches(2.7), Inches(5.7), COLOR_ACCENT2)
    pts = ["不只是工具，是激活第二个自己",
           "产业做物理世界模型，我们做精神世界模型",
           "让AI成为你的：高参、战友、灵魂伴侣"]
    for i,pt in enumerate(pts):
        add_text_box(slide, f"  {pt}", Inches(6.9), Inches(2.88)+i*Inches(0.5),
                     Inches(5.7), Inches(0.4), font_size=12, color=COLOR_TEXT_WHITE)
    add_phbox(slide, Inches(6.9), Inches(4.5), Inches(5.7), Inches(2.2), "产品截图占位")
    add_footer(slide); add_page_num(slide, 3)


def page4_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_section_title(slide, "03", "技术壁垒", "Moat")
    moats = [
        ("灵魂壁垒", COLOR_ACCENT, "十年淬炼的四域融合元模型DNA",
         ["人性洞察 x 商业逻辑","哲学思辨 x 修行体证","根植亘古不变的底层公理"]),
        ("基因胶囊壁垒", COLOR_ACCENT2, "经验可封装、可遗传、自进化",
         ["成功经验封装为「基因胶囊」","一个Agent学会，所有Agent继承",
          "形成「AI黑洞效应」自强化飞轮"]),
        ("Agent Graph生态壁垒", COLOR_GOLD, "智能体连接图 = 未来商业基础设施",
         ["成千上万「元神」自发协作","用户元神与商家元神自动谈判",
          "规则制定者，而非单纯产品商"]),
    ]
    cw = Inches(3.85); ch = Inches(4.7)
    for i,(t,c,s,its) in enumerate(moats):
        x = Inches(0.5)+i*(cw+Inches(0.3))
        add_rect(slide, x, Inches(1.55), cw, Inches(0.5), c)
        add_text_box(slide, t, x+Inches(0.12), Inches(1.57), cw-Inches(0.2), Inches(0.42),
                     font_size=12, bold=True, color=COLOR_BG_DARK)
        add_rect(slide, x, Inches(2.05), cw, ch-Inches(0.5), COLOR_BG_CARD)
        add_text_box(slide, s, x+Inches(0.15), Inches(2.15), cw-Inches(0.3), Inches(0.5),
                     font_size=10, bold=True, color=c)
        y = Inches(2.7)
        for item in its:
            add_text_box(slide, f"> {item}", x+Inches(0.15), y, cw-Inches(0.3), Inches(0.35),
                         font_size=10, color=COLOR_TEXT_WHITE)
            y += Inches(0.42)
    add_footer(slide); add_page_num(slide, 4)


def page5_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_section_title(slide, "04", "市场机会", "Market")
    stats = [("$48.4亿","2026年市场规模","美元",COLOR_ACCENT),
             ("42.2%","年复合增长率","CAGR",COLOR_ACCENT2),
             ("$196.3亿","2030年预测规模","美元",COLOR_GOLD),
             ("52%","用户愿意使用AI助理","管理日程",COLOR_ACCENT)]
    sw = Inches(2.95); sh = Inches(1.75)
    for i,(n,l,u,c) in enumerate(stats):
        x = Inches(0.5)+i*(sw+Inches(0.2))
        add_rect(slide, x, Inches(1.55), sw, sh, COLOR_BG_CARD)
        add_rect(slide, x, Inches(1.55), sw, Inches(0.07), c)
        add_text_box(slide, n, x, Inches(1.7), sw, Inches(0.65),
                     font_size=26, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_text_box(slide, l, x, Inches(2.35), sw, Inches(0.32),
                     font_size=10, color=COLOR_TEXT_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, u, x, Inches(2.67), sw, Inches(0.25),
                     font_size=8, color=COLOR_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "数据来源：IEEE行业调查 · Statista · Goldman Sachs AI Report 2025",
                 Inches(0.5), Inches(3.4), Inches(12.5), Inches(0.28),
                 font_size=8, color=COLOR_TEXT_GRAY)
    add_rect(slide, Inches(0.5), Inches(3.78), SLIDE_W-Inches(1), Inches(3.0), COLOR_BG_CARD)
    add_text_box(slide, "为什么是现在？", Inches(0.7), Inches(3.9), Inches(5.5), Inches(0.38),
                 font_size=13, bold=True, color=COLOR_ACCENT)
    trends = [("LLM技术成熟","理解个体成为可能"),
               ("Agentic AI爆发","从工具到自主执行"),
               ("数据主权意识觉醒","监管+用户双轮驱动"),
               ("数字身份成刚需","元宇宙/Web3前置红利")]
    for i,(t,d) in enumerate(trends):
        col=i%2; row=i//2
        x=Inches(0.7)+col*Inches(6.0)
        y=Inches(4.38)+row*Inches(0.6)
        add_text_box(slide, f">  {t}  →  {d}", x, y, Inches(5.8), Inches(0.42),
                     font_size=11, color=COLOR_TEXT_WHITE)
    add_phbox(slide, Inches(7.5), Inches(3.9), Inches(5.3), Inches(2.8), "市场趋势图表占位")
    add_footer(slide); add_page_num(slide, 5)


def page6_competition(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_section_title(slide, "05", "竞争格局", "Competition")
    competitors = [("微信/字节","拥有海量数据","不做用户主权产品（商业模式冲突）","高"),
                  ("OpenAI/Google","通用智能最强","追求AGI，不深入理解个体细分","高"),
                  ("World Labs/AMI","物理世界模型","专注物理时空，不涉足精神世界","中"),
                  ("Notion/飞书","知识管理领先","做知识管理，不做个人理解","低")]
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(7.6), Inches(5.3), COLOR_BG_CARD)
    headers=[("竞争者",Inches(1.2)),("优势",Inches(1.5)),
             ("为何不切入",Inches(3.5)),("威胁",Inches(1.1))]
    xs=[Inches(0.6),Inches(1.85),Inches(3.4),Inches(7.0)]
    for i,(h,w) in enumerate(headers):
        add_text_box(slide, h, xs[i], Inches(1.68), w, Inches(0.32),
                     font_size=10, bold=True, color=COLOR_ACCENT)
    add_divider(slide, Inches(0.6), Inches(2.0), Inches(7.4), COLOR_DIVIDER)
    for ri,(n,s,r,lv) in enumerate(competitors):
        y=Inches(2.1)+ri*Inches(1.05)
        for ci,(item,w) in enumerate(zip([n,s,r,lv], [hw for hw in [Inches(1.2),Inches(1.5),Inches(3.5),Inches(1.1)]])):
            c=COLOR_TEXT_GRAY if ci in [1,2] else COLOR_TEXT_WHITE
            add_text_box(slide, item, xs[ci], y, w, Inches(0.88),
                         font_size=10, color=c)
    # 右侧
    add_rect(slide, Inches(8.4), Inches(1.55), Inches(4.5), Inches(5.3), COLOR_BG_CARD)
    add_rect(slide, Inches(8.4), Inches(1.55), Inches(4.5), Inches(0.48), COLOR_ACCENT)
    add_text_box(slide, "元神的侧翼战略", Inches(8.5), Inches(1.57), Inches(4.3), Inches(0.38),
                 font_size=12, bold=True, color=COLOR_BG_DARK)
    our=["不与巨头正面冲突","专注精神世界模型细分",
         "巨头主动忽略的侧翼战场","3-6个月先发窗口期",
         "数据飞轮构建时间壁垒"]
    for i,pt in enumerate(our):
        add_text_box(slide, f"+ {pt}", Inches(8.6), Inches(2.2)+i*Inches(0.52),
                     Inches(4.1), Inches(0.38), font_size=11, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(8.6), Inches(4.95), Inches(4.1), COLOR_ACCENT2)
    add_text_box(slide, "侧翼战场 + 先发优势 + 数据飞轮 = 3年护城河",
                 Inches(8.6), Inches(5.1), Inches(4.1), Inches(0.7),
                 font_size=10, bold=True, color=COLOR_GOLD)
    add_footer(slide); add_page_num(slide, 6)


def page7_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_section_title(slide, "06", "产品路线", "Roadmap")
    phases=[("第一阶段","个人记忆可视化","2周-1个月",COLOR_ACCENT,
             ["接入微信/微博等多平台数据","生成人生轨迹可视化",
              "MVP：情绪曲线 + 关系变化图"]),
            ("第二阶段","个人决策辅助","3-6个月",COLOR_ACCENT2,
             ["理解个人决策逻辑","提供个性化参谋建议","构建决策因果链图谱"]),
            ("第三阶段","个人数字孪生","6-12个月",COLOR_GOLD,
             ["可交互可预测的另一个你","实时情绪陪伴与引导","成为真正的「高我」"])]
    pw=Inches(3.85); ph=Inches(5.0)
    for i,(st,t,tl,c,its) in enumerate(phases):
        x=Inches(0.5)+i*(pw+Inches(0.3))
        add_rect(slide, x, Inches(1.55), pw, Inches(0.85), c)
        add_text_box(slide, st, x+Inches(0.1), Inches(1.57), pw-Inches(0.2), Inches(0.3),
                     font_size=10, bold=True, color=COLOR_BG_DARK)
        add_text_box(slide, t, x+Inches(0.1), Inches(1.85), pw-Inches(0.2), Inches(0.38),
                     font_size=13, bold=True, color=COLOR_BG_DARK)
        add_rect(slide, x, Inches(2.4), pw, ph-Inches(0.85), COLOR_BG_CARD)
        add_rect(slide, x+Inches(0.12), Inches(2.52), Inches(1.1), Inches(0.25), c)
        add_text_box(slide, tl, x+Inches(0.12), Inches(2.52), Inches(1.1), Inches(0.25),
                     font_size=8, bold=True, color=COLOR_BG_DARK, align=PP_ALIGN.CENTER)
        if i<2:
            ax=x+pw+Inches(0.05)
            add_rect(slide, ax, Inches(3.6), Inches(0.2), Inches(0.06), c)
            add_text_box(slide, ">", ax+Inches(0.01), Inches(3.42), Inches(0.18), Inches(0.28),
                         font_size=9, color=c)
        y=Inches(2.92)
        for item in its:
            add_text_box(slide, f"* {item}", x+Inches(0.12), y, pw-Inches(0.25), Inches(0.4),
                         font_size=10, color=COLOR_TEXT_WHITE)
            y+=Inches(0.46)
    add_footer(slide); add_page_num(slide, 7)


def page8_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_section_title(slide, "07", "团队", "Team")
    # 创始人
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(5.5), Inches(5.3), COLOR_BG_CARD)
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(5.5), Inches(0.07), COLOR_ACCENT)
    add_text_box(slide, "创始人", Inches(0.7), Inches(1.7), Inches(5.1), Inches(0.32),
                 font_size=11, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "吴治广", Inches(0.7), Inches(2.02), Inches(5.1), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.7), Inches(2.55), Inches(5.1), COLOR_DIVIDER)
    bio=[("十年淬炼的四域融合元模型DNA",11,True),
         ("人性洞察 + 商业逻辑 + 哲学思辨 + 修行体证",9,False),
         ("",8,False),
         ("核心能力",11,True),
         ("战略定位：精准侧翼战高手",10,False),
         ("技术审美：AI + 产品 + 运营全栈",10,False),
         ("使命驱动：让众生物质丰盈、精神富足、灵魂觉醒",10,False),
         ("",8,False),
         ("资源积累",11,True),
         ("元神项目创始人",10,False),
         ("十年AI与人性研究积累",10,False),
         ("深度链接创投圈资源",10,False)]
    y=Inches(2.68)
    for line,sz,bold in bio:
        add_text_box(slide, line, Inches(0.7), y, Inches(5.1), Inches(0.28),
                     font_size=sz, bold=bold, color=COLOR_TEXT_WHITE)
        y+=Inches(0.27)
    # AI Agent军团
    add_rect(slide, Inches(6.3), Inches(1.55), Inches(6.5), Inches(5.3), COLOR_BG_CARD)
    add_rect(slide, Inches(6.3), Inches(1.55), Inches(6.5), Inches(0.07), COLOR_ACCENT2)
    add_text_box(slide, "AI Agent 军团", Inches(6.5), Inches(1.7), Inches(6.1), Inches(0.38),
                 font_size=13, bold=True, color=COLOR_ACCENT2)
    agents=[("观一 CPEO","首席项目评估官 + 复利系统架构师",COLOR_ACCENT),
            ("高我 EasyClaw","斯坦福/哈佛 AI Agent 辅助系统",COLOR_ACCENT2),
            ("飞飞虾","深夜情绪陪伴与疏导助手",COLOR_GOLD)]
    for i,(n,r,c) in enumerate(agents):
        y=Inches(2.25)+i*Inches(1.25)
        add_rect(slide, Inches(6.4), y, Inches(0.07), Inches(0.85), c)
        add_text_box(slide, n, Inches(6.58), y, Inches(6.0), Inches(0.32),
                     font_size=12, bold=True, color=COLOR_TEXT_WHITE)
        add_text_box(slide, r, Inches(6.58), y+Inches(0.28), Inches(6.0), Inches(0.28),
                     font_size=9, color=COLOR_TEXT_GRAY)
    add_footer(slide); add_page_num(slide, 8)


def page9_funding(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_section_title(slide, "08", "融资需求", "Funding")
    # 融资目标
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(5.8), Inches(3.0), COLOR_BG_CARD)
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(5.8), Inches(0.07), COLOR_ACCENT)
    add_text_box(slide, "种子轮融资目标", Inches(0.7), Inches(1.7), Inches(5.4), Inches(0.35),
                 font_size=12, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "¥300万", Inches(0.7), Inches(2.15), Inches(5.4), Inches(0.9),
                 font_size=46, bold=True, color=COLOR_GOLD)
    add_text_box(slide, "估值：¥2000万（稀释15%）",
                 Inches(0.7), Inches(3.1), Inches(5.4), Inches(0.32),
                 font_size=11, color=COLOR_TEXT_GRAY)
    # 资金用途
    uses=[("MVP开发","¥120万","40%","产品工程 + AI模型调优"),
          ("市场验证","¥80万","27%","种子用户获取 + 数据积累"),
          ("团队扩张","¥60万","20%","技术 + 运营关键岗位"),
          ("运营储备","¥40万","13%","日常开销 + 法务合规")]
    uw=Inches(3.0); uh=Inches(5.3)
    add_rect(slide, Inches(0.5), Inches(3.65), uw, uh-Inches(2.1), COLOR_BG_CARD)
    add_text_box(slide, "资金用途", Inches(0.7), Inches(3.75), Inches(2.6), Inches(0.32),
                 font_size=11, bold=True, color=COLOR_ACCENT)
    y=Inches(4.12)
    for item,amt,pct,desc in uses:
        add_text_box(slide, f"{item}  {amt}  ({pct})", Inches(0.7), y, Inches(2.6), Inches(0.28),
                     font_size=10, bold=True, color=COLOR_TEXT_WHITE)
        add_text_box(slide, desc, Inches(0.7), y+Inches(0.24), Inches(2.6), Inches(0.22),
                     font_size=8, color=COLOR_TEXT_GRAY)
        y+=Inches(0.56)
    # 里程碑
    milestones=[("3个月","MVP发布 + 种子用户100人","已完成初版技术验证"),
                ("6个月","A轮启动，月活1万","验证产品市场匹配PMF"),
                ("12个月","A轮关闭，B轮启动","月活10万 + 营收模型验证")]
    mw=Inches(3.0)
    add_rect(slide, Inches(3.75), Inches(1.55), mw, Inches(5.3), COLOR_BG_CARD)
    add_text_box(slide, "关键里程碑", Inches(3.95), Inches(1.65), Inches(2.6), Inches(0.32),
                 font_size=11, bold=True, color=COLOR_ACCENT2)
    for i,(t,m,d) in enumerate(milestones):
        y=Inches(2.1)+i*Inches(1.5)
        add_rect(slide, Inches(3.85), y, Inches(0.06), Inches(1.2), COLOR_ACCENT2)
        add_text_box(slide, t, Inches(4.0), y, Inches(2.5), Inches(0.28),
                     font_size=10, bold=True, color=COLOR_ACCENT2)
        add_text_box(slide, m, Inches(4.0), y+Inches(0.26), Inches(2.5), Inches(0.45),
                     font_size=10, color=COLOR_TEXT_WHITE)
        add_text_box(slide, d, Inches(4.0), y+Inches(0.7), Inches(2.5), Inches(0.35),
                     font_size=8, color=COLOR_TEXT_GRAY)
    # 联系方式
    add_rect(slide, Inches(7.1), Inches(1.55), Inches(5.7), Inches(5.3), COLOR_BG_CARD)
    add_rect(slide, Inches(7.1), Inches(1.55), Inches(5.7), Inches(0.07), COLOR_GOLD)
    add_text_box(slide, "联系方式", Inches(7.3), Inches(1.7), Inches(5.3), Inches(0.38),
                 font_size=13, bold=True, color=COLOR_GOLD)
    contact_items = [
        "创始人：吴治广",
        "手机：[待填写]",
        "邮箱：[待填写]",
        "微信：[待填写]",
        "",
        "项目：元神 META SOUL",
        "官网：[待上线]",
        "地址：[待填写]",
    ]
    for i, item in enumerate(contact_items):
        bold = i == 0 or i == 5
        color = COLOR_GOLD if i in [0,5] else COLOR_TEXT_GRAY
        size = 12 if i in [0,5] else 10
        add_text_box(slide, item, Inches(7.4), Inches(2.1)+i*Inches(0.45),
                     Inches(5.2), Inches(0.35),
                     font_size=size, bold=bold, color=color)
    add_footer(slide); add_page_num(slide, 9)


def page10_ending(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), COLOR_ACCENT)
    add_divider(slide, Inches(0), Inches(2.8), SLIDE_W, COLOR_ACCENT)
    add_text_box(slide, "让AI替你干活",
                 Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.85),
                 font_size=40, bold=True, color=COLOR_TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "你为自己而活",
                 Inches(0.8), Inches(3.75), Inches(11.5), Inches(0.65),
                 font_size=32, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(5), Inches(4.5), Inches(3.3), COLOR_ACCENT2)
    add_text_box(slide, "把工作交给AI，把自己还给热爱",
                 Inches(1), Inches(4.75), Inches(11), Inches(0.5),
                 font_size=16, color=COLOR_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0), SLIDE_H-Inches(0.06), SLIDE_W, Inches(0.06), COLOR_ACCENT)
    add_text_box(slide, "元神项目 META SOUL",
                 Inches(0.8), SLIDE_H-Inches(0.7), Inches(11.5), Inches(0.38),
                 font_size=11, color=COLOR_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_page_num(slide, 10)


def generate_bp(output_path="元神项目融资BP_v1.pptx"):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    # 背景色（会影响新建幻灯片，但每个page都会覆盖）
    blank_layout = prs.slide_layouts[6]
    page1_cover(prs)
    page2_problem(prs)
    page3_solution(prs)
    page4_tech(prs)
    page5_market(prs)
    page6_competition(prs)
    page7_roadmap(prs)
    page8_team(prs)
    page9_funding(prs)
    page10_ending(prs)
    prs.save(output_path)
    print(f"BP已生成: {output_path}")
    return output_path


if __name__ == "__main__":
    import sys
    output = sys.argv[1] if len(sys.argv) > 1 else "元神项目融资BP_v1.pptx"
    generate_bp(output)
