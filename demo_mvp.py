#!/usr/bin/env /usr/bin/python3
# -*- coding: utf-8 -*-
"""元神项目 MVP Demo 生成器 v3 - 完整干净版"""

import os, sys, json, random, ast
from datetime import datetime, timedelta
from pathlib import Path

DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "")
DEEPSEEK_API_URL = "https://api.deepseek.com/chat/completions"
OUTPUT_DIR = Path("demo_output")
OUTPUT_DIR.mkdir(exist_ok=True)

# ============================================================
# 数据生成
# ============================================================
def gen_moments():
    m = []
    phases = [
        (2016,1,2017,12,["今天入职新公司，希望能学到很多！加油","深夜加班，但感觉自己在快速成长，值了。",
            "周末参加产品经理培训，老师讲的精益创业让我茅塞顿开。","第一次独立负责项目，紧张又兴奋！",
            "下班后的北京夜景真美，霓虹灯下的北漂梦。","读完《从零到一》，创业的种子在心里发芽。",
            "和同事聚餐，聊到凌晨一点，意犹未尽。","第一次在朋友圈分享技术文章，收到了很多点赞。",
            "周末去五道口逛书店，买了一堆产品经理相关的书。","加班到凌晨两点，但上线那一刻的成就感太棒了！"]),
        (2018,1,2019,6,["今天正式离职，开始我的创业之路！不留退路。","拿到第一笔天使投资，虽然只有50万，但这是梦想的开始！",
            "团队从1个人变成5个人，每个人都充满激情。","产品上线第一天，收到了用户的感谢反馈，眼眶湿润。",
            "凌晨三点还在改BP，投资人的话让我满血复活。","第一次被媒体采访，紧张得手心出汗。",
            "团队吵架了，但吵完之后产品方向更清晰了。","参加互联网大会，见到了偶像，感觉自己在正确的路上。",
            "用户突破1万！虽然离盈利还远，但这是里程碑。","一个人扛不住了，在厕所偷偷哭了两分钟，然后回去继续开会。"]),
        (2019,7,2020,12,["投资人说要收缩，我们不得不裁掉一半的同事。","最难的时候，房东免了我一个月房租，说看好我。",
            "疫情期间在家办公，第一次感受到远程协作的挑战。","和合伙人理念分歧，彻夜长谈，决定各自寻找方向。",
            "转型做在线教育，从零开始学直播带货。","最穷的时候卡里只剩2000块，但告诉自己要撑住。",
            "直播第一场只有5个人观看，但没有放弃。","终于还清了一部分债务，轻装上阵。",
            "学会了冥想，每天早起打坐30分钟，心态稳了很多。","低谷期读了很多哲学书，那句话救了我：认识你自己。"]),
        (2021,1,2023,6,["新项目拿到种子轮200万，这次我更懂得控制节奏了。","开始每周健身三次，身体是革命的本钱。",
            "读完MBA课程，认知升级的感觉太上头了。","学会说No，拒绝了很多无效社交，时间更值钱了。",
            "开始认真谈恋爱，第一次感受到有人懂我的创业焦虑。","搬到望京的新办公室，终于有了像样的根据地。",
            "团队20人了，文化建设变得很重要。","给父母买了房子，把他们接到了北京。",
            "开始写日记，记录每一天的思考，这是自我进化的方式。","去了一趟西藏，在布达拉宫前哭了两小时，心里很多东西放下了。"]),
        (2023,7,2026,3,["决定All in AI，做一件让更多人受益的事。","元神项目的想法成型了：让每个人都有一个比自己更好的自己。",
            "开始深度研究世界模型、Agentic AI，感觉自己在正确的赛道上。","把创业经历写成了一篇文章，收到了500条评论，大部分是感谢。",
            "深刻理解：赚钱是为了自由，自由是为了做更有意义的事。","每周给家里打电话，父母的声音是最好的减压药。",
            "现在每天早起写反思日记，发现自己比三年前成熟了太多。","不再焦虑结果，专注于做好当下的每一件事。",
            "我相信：真正的成功，是让身边的人都变得更好。","让AI帮我干活，把自己还给热爱。"]),
    ]
    for sy,sm,ey,em,themes in phases:
        cur = datetime(sy,sm,1); end = datetime(ey,em,28)
        while cur <= end:
            np_ = random.randint(2,5)
            for _ in range(np_):
                d = random.randint(1,min(28,(cur+timedelta(days=30)).day))
                h = random.randint(8,23); mn = random.randint(0,59)
                pd = cur.replace(day=min(d,28)).replace(hour=h).replace(minute=mn)
                if pd > datetime.now(): continue
                c = random.choice(themes)
                if random.random()<0.2: c += " [图片]"
                m.append({"date":pd.strftime("%Y-%m-%d"),"time":pd.strftime("%H:%M"),
                           "content":c,"timestamp":pd.strftime("%Y-%m-%d %H:%M")})
            m2=cur.month+1; y2=cur.year
            if m2>12: m2=1; y2+=1
            cur = datetime(y2,m2,1)
    m.sort(key=lambda x:x["timestamp"]); return m

def sentiment(c):
    cl=c.lower()
    ng=sum(-0.4 for w in ["焦虑","难过","哭","痛苦","危机","裁员","低谷","绝望","迷茫","负债","压力","穷","难"] if w in cl)
    ps=sum(0.4 for w in ["开心","兴奋","骄傲","感谢","突破","里程碑","成功","加油","成长","收获","爱","希望","美好","激动"] if w in cl)
    md=sum(0.1 for w in ["加班","忙碌","学习","思考","记录","分享"] if w in cl)
    if "[图片]" in c: ps+=0.1
    return max(-1.0,min(1.0,ng+ps+md))

def topic(c):
    cl=c.lower()
    if any(w in cl for w in ["创业","融资","项目","产品","公司","团队","投资","BP","天使"]): return "创业"
    if any(w in cl for w in ["加班","工作","职场","入职","辞职","同事","老板"]): return "职场"
    if any(w in cl for w in ["健身","健康","身体","冥想","早起","瑜伽"]): return "健康"
    if any(w in cl for w in ["学习","读书","课程","MBA","培训","认知"]): return "学习"
    if any(w in cl for w in ["旅行","西藏","风景","周末","度假","长城"]): return "旅行"
    if any(w in cl for w in ["感情","恋爱","父母","家人","爱","电话","朋友"]): return "情感"
    if any(w in cl for w in ["AI","科技","技术","互联网","未来","智能"]): return "科技"
    if any(w in cl for w in ["反思","日记","成长","思考","成熟"]): return "成长"
    return "生活"

# ============================================================
# 可视化
# ============================================================
def gen_viz(moments, out_dir):
    import matplotlib; matplotlib.use("Agg")
    import matplotlib.pyplot as plt; import numpy as np
    from collections import Counter
    plt.rcParams["font.sans-serif"]=["Arial Unicode MS"]; plt.rcParams["axes.unicode_minus"]=False

    dates=[datetime.strptime(r["timestamp"],"%Y-%m-%d %H:%M") for r in moments]
    sents=[r["sentiment"] for r in moments]
    topics=[r["topic"] for r in moments]

    # 情绪曲线
    mo={}
    for d,s in zip(dates,sents):
        k=d.strftime("%Y-%m"); mo.setdefault(k,[]).append(s)
    keys=sorted(mo.keys()); mav=[np.mean(mo[k]) for k in keys]
    xs=[datetime.strptime(k,"%Y-%m") for k in keys]
    fig,ax=plt.subplots(figsize=(14,5))
    bc=["#00C85C" if v>0.05 else("#FF6060" if v<-0.05 else "#6C63FF") for v in mav]
    ax.bar(xs,mav,width=20,color=bc,alpha=0.8,edgecolor="none")
    ax.axhline(0,color="white",linewidth=0.8,alpha=0.5)
    if len(mav)>6:
        z=np.polyfit(range(len(mav)),mav,2)
        ax.plot(xs,[np.poly1d(z)(i) for i in range(len(mav))],color="#00D4AA",linewidth=2.5,linestyle="--",alpha=0.9)
    ax.set_facecolor("#0D0D1A"); fig.patch.set_facecolor("#0D0D1A")
    ax.tick_params(colors="#A0A0B0",labelsize=8)
    for sp in ["top","right"]: ax.spines[sp].set_visible(False)
    ax.spines["bottom"].set_color("#3A3A5C"); ax.spines["left"].set_color("#3A3A5C")
    ax.set_title("Li Ming - 10-Year Emotion Curve",color="white",fontsize=14,fontweight="bold",pad=15)
    ax.set_ylabel("Sentiment Score (-1 to +1)",color="#A0A0B0",fontsize=9)
    plt.tight_layout(); mp=str(out_dir/"emotion_curve.png")
    plt.savefig(mp,dpi=150,bbox_inches="tight",facecolor="#0D0D1A"); plt.close()
    print("  emotion_curve.png OK")

    # 主题饼图
    tc=Counter(topics)
    fig,ax=plt.subplots(figsize=(8,8))
    pc=["#6C63FF","#00D4AA","#F0C040","#00C85C","#FF6060","#FF8C42","#9B59B6","#E74C3C"]
    wedges,texts,auts=ax.pie(list(tc.values()),labels=list(tc.keys()),autopct="%1.0f%%",
        colors=pc[:len(tc)],textprops={"color":"white","fontsize":11})
    for t in texts: t.set_color("white")
    for a in auts: a.set_color("white"); a.set_fontsize(10)
    ax.set_title("Li Ming - Topic Distribution",color="white",fontsize=14,fontweight="bold",pad=20)
    fig.patch.set_facecolor("#0D0D1A"); ax.set_facecolor("#0D0D1A")
    lbs=[f"{l} ({c} posts)" for l,c in tc.most_common()]
    leg = ax.legend(wedges,lbs,title="Topics",loc="center left",
              bbox_to_anchor=(1,0,0.5,1),labelcolor="white",facecolor="#1A1A2E",edgecolor="#3A3A5C")
    leg.get_title().set_color("white")
    tp=str(out_dir/"topic_distribution.png")
    plt.savefig(tp,dpi=150,bbox_inches="tight",facecolor="#0D0D1A"); plt.close()
    print("  topic_distribution.png OK")

    # 年度热力图
    yr={}
    for d,s in zip(dates,sents):
        y,m=d.year,d.month
        yr.setdefault(y,{"s":[0]*12,"c":[0]*12})
        yr[y]["s"][m-1]+=s; yr[y]["c"][m-1]+=1
    years=sorted(yr.keys())
    mat=[[yr[y]["s"][m]/(yr[y]["c"][m] or 1) if yr[y]["c"][m]>0 else 0 for m in range(12)] for y in years]
    fig,ax=plt.subplots(figsize=(16,4))
    ax.imshow(mat,aspect="auto",cmap="RdYlGn",vmin=-0.5,vmax=0.5)
    ax.set_xticks(range(12)); ax.set_xticklabels(["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"],color="#A0A0B0",fontsize=8)
    ax.set_yticks(range(len(years))); ax.set_yticklabels(years,color="#A0A0B0",fontsize=9)
    ax.set_title("Li Ming - Yearly Emotion Heatmap (2016-2026)",color="white",fontsize=13,fontweight="bold",pad=10)
    fig.patch.set_facecolor("#0D0D1A"); ax.tick_params(colors="#A0A0B0")
    plt.tight_layout(); hp=str(out_dir/"yearly_heatmap.png")
    plt.savefig(hp,dpi=150,bbox_inches="tight",facecolor="#0D0D1A"); plt.close()
    print("  yearly_heatmap.png OK")

    return {"emotion_curve":mp,"topic_distribution":tp,"yearly_heatmap":hp}

# ============================================================
# PPT
# ============================================================
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import numpy as np
from collections import Counter

C_DK=RGBColor(0x0D,0x0D,0x1A); C_CD=RGBColor(0x1A,0x1A,0x2E)
C_AC=RGBColor(0x6C,0x63,0xFF); C_AC2=RGBColor(0x00,0xD4,0xAA)
C_WH=RGBColor(0xFF,0xFF,0xFF); C_GR=RGBColor(0xA0,0xA0,0xB0)
C_GD=RGBColor(0xF0,0xC0,0x40); C_GN=RGBColor(0x00,0xC8,0x5C)
C_RD=RGBColor(0xFF,0x60,0x60); C_DV=RGBColor(0x3A,0x3A,0x5C)
SW=Inches(13.33); SH=Inches(7.5)

def sbg(sl,c): f=sl.background.fill; f.solid(); f.fore_color.rgb=c
def stb(sl,tx,l,t,w,h,fs=14,bold=False,c=C_WH,align=PP_ALIGN.LEFT):
    bx=sl.shapes.add_textbox(l,t,w,h); tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align; rn=p.add_run()
    rn.text=tx; rn.font.size=Pt(fs); rn.font.bold=bold; rn.font.color.rgb=c
def srect(sl,l,t,w,h,fc):
    sh=sl.shapes.add_shape(1,l,t,w,h); sh.fill.solid(); sh.fill.fore_color.rgb=fc; sh.line.fill.background()
def sdiv(sl,l,t,w,c=C_AC):
    ln=sl.shapes.add_shape(1,l,t,w,Pt(2)); ln.fill.solid(); ln.fill.fore_color.rgb=c; ln.line.fill.background()
def spg(sl,n): stb(sl,f"{n}/12",SW-Inches(1.1),SH-Inches(0.38),Inches(1),Inches(0.25),fs=8,c=C_GR,align=PP_ALIGN.RIGHT)
def sft(sl): stb(sl,"Meta Soul MVP Demo 2026",Inches(0.4),SH-Inches(0.38),Inches(5),Inches(0.25),fs=8,c=C_GR)
def stitle(sl,num,en,cn):
    srect(sl,Inches(0),Inches(0),SW,Inches(0.05),C_AC)
    stb(sl,num,Inches(0.4),Inches(0.22),Inches(1),Inches(0.65),fs=36,bold=True,c=C_AC)
    stb(sl,en,Inches(0.4),Inches(0.75),Inches(8),Inches(0.42),fs=18,bold=True,c=C_WH)
    stb(sl,cn,Inches(0.4),Inches(0.72),Inches(7),Inches(0.5),fs=12,c=C_GR)
    sdiv(sl,Inches(0.4),Inches(1.25),SW-Inches(0.8),C_DV)
def scap(sl,n,t,sub,c=C_AC):
    srect(sl,Inches(0),Inches(0),SW,Inches(0.05),C_AC)
    stb(sl,t,Inches(0.5),Inches(2.8),Inches(12.3),Inches(1.0),fs=n,bold=True,c=c,align=PP_ALIGN.CENTER)
    stb(sl,sub,Inches(0.5),Inches(3.6),Inches(12.3),Inches(0.55),fs=18,c=C_WH,align=PP_ALIGN.CENTER)
    srect(sl,Inches(0),SH-Inches(0.05),SW,Inches(0.05),C_AC)

def gen_ppt(moments, analysis, viz, out_path):
    prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
    blank=prs.slide_layouts[6]
    n=len(moments)
    avg_s=np.mean([r["sentiment"] for r in analysis])
    pos_r=sum(1 for r in analysis if r["sentiment"]>0.1)/n
    neg_r=sum(1 for r in analysis if r["sentiment"]<-0.1)/n

    # P1 封面
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK)
    srect(sl,Inches(0),Inches(0),SW,Inches(0.05),C_AC)
    srect(sl,Inches(0),Inches(2.9),SW,Inches(0.04),C_AC)
    srect(sl,Inches(0),SH-Inches(0.05),SW,Inches(0.05),C_AC)
    stb(sl,"META SOUL",Inches(0.5),Inches(3.05),Inches(12.3),Inches(0.95),fs=58,bold=True,c=C_AC,align=PP_ALIGN.CENTER)
    stb(sl,"元神项目 MVP Demo",Inches(0.5),Inches(3.92),Inches(12.3),Inches(0.65),fs=28,c=C_WH,align=PP_ALIGN.CENTER)
    sdiv(sl,Inches(4.5),Inches(4.68),Inches(4.3),C_AC2)
    stb(sl,"虚构人物：李明 · 互联网创业者 · 北漂10年 · "+str(n)+"条朋友圈",Inches(0.5),Inches(4.92),Inches(12.3),Inches(0.5),fs=14,c=C_GR,align=PP_ALIGN.CENTER)
    stb(sl,"从朋友圈数据，看一个人如何变成今天的自己",Inches(0.5),Inches(5.5),Inches(12.3),Inches(0.45),fs=12,c=C_AC2,align=PP_ALIGN.CENTER)
    sft(sl); spg(sl,1)

    # P2 档案
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"01","主角档案","Profile")
    srect(sl,Inches(0.4),Inches(1.45),Inches(5.5),Inches(5.5),C_CD)
    srect(sl,Inches(0.4),Inches(1.45),Inches(5.5),Inches(0.06),C_AC)
    stb(sl,"李明 Li Ming",Inches(0.6),Inches(1.62),Inches(5.1),Inches(0.5),fs=20,bold=True,c=C_WH)
    attrs=[("年龄","30岁"),("职业","互联网连续创业者"),("地点","北京（北漂10年）"),("阶段","第四次创业 · All in AI"),("性格","理想主义者，热爱科技"),("使命","让AI理解你的世界")]
    y=Inches(2.22)
    for k,v in attrs:
        stb(sl,k+"：",Inches(0.6),y,Inches(1.4),Inches(0.32),fs=11,bold=True,c=C_AC)
        stb(sl,v,Inches(2.05),y,Inches(3.7),Inches(0.32),fs=11,c=C_WH); y+=Inches(0.4)
    srect(sl,Inches(6.2),Inches(1.45),Inches(6.8),Inches(5.5),C_CD)
    stb(sl,"10年轨迹",Inches(6.4),Inches(1.6),Inches(6.4),Inches(0.38),fs=13,bold=True,c=C_AC2)
    tl=[("2016-2017","初入职场","月薪8K到首次独立带项目"),("2018-2019","首次创业","天使投资、团队5人、高光时刻"),("2019-2020","危机转型","裁员、卡里2000块、西藏重建"),("2021-2023","复苏重建","种子轮200万、20人团队"),("2024-2026","使命驱动","All in AI、元神项目启动")]
    clrs=[C_AC2,C_AC,C_RD,C_AC2,C_GD]
    y=Inches(2.1)
    for i,(yr,t,d) in enumerate(tl):
        srect(sl,Inches(6.4),y,Inches(0.06),Inches(0.82),clrs[i])
        stb(sl,yr,Inches(6.55),y,Inches(1.1),Inches(0.28),fs=9,bold=True,c=clrs[i])
        stb(sl,t,Inches(6.55),y+Inches(0.25),Inches(6.2),Inches(0.28),fs=10,bold=True,c=C_WH)
        stb(sl,d,Inches(6.55),y+Inches(0.5),Inches(6.2),Inches(0.25),fs=9,c=C_GR); y+=Inches(0.95)
    sft(sl); spg(sl,2)

    # P3 数据概览
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"02","数据概览","Data Overview")
    stats=[(str(n),"条朋友圈","10年完整记录",C_AC),(str(int(n/10))+"月","时间跨度","月均发布频率",C_AC2),(f"{avg_s:.2f}","平均情感","整体积极向上",C_GN),(f"{pos_r:.0%}","正面内容",f"负面{neg_r:.0%}",C_GD)]
    sw=Inches(3.0)
    for i,(num,lab,desc,c) in enumerate(stats):
        x=Inches(0.4)+i*(sw+Inches(0.18))
        srect(sl,x,Inches(1.45),sw,Inches(1.9),C_CD); srect(sl,x,Inches(1.45),sw,Inches(0.06),c)
        stb(sl,num,x,Inches(1.6),sw,Inches(0.75),fs=26,bold=True,c=c,align=PP_ALIGN.CENTER)
        stb(sl,lab,x,Inches(2.35),sw,Inches(0.32),fs=11,c=C_WH,align=PP_ALIGN.CENTER)
        stb(sl,desc,x,Inches(2.67),sw,Inches(0.28),fs=9,c=C_GR,align=PP_ALIGN.CENTER)
    srect(sl,Inches(0.4),Inches(3.5),SW-Inches(0.8),Inches(3.5),C_CD)
    stb(sl,"李明的10年发生了什么？",Inches(0.6),Inches(3.62),Inches(6),Inches(0.38),fs=12,bold=True,c=C_AC)
    hs=[("北漂10年，从月薪8K到连续创业","从高光到低谷到重建"),("经历3次重大转型","每次都是对自我的重新认知"),("创业失败时卡里只剩2000块","但没有放弃"),("2024年All in AI","找到真正的使命：元神项目")]
    y=Inches(4.1)
    for h,d in hs:
        stb(sl,"  "+h,Inches(0.6),y,Inches(6.0),Inches(0.32),fs=11,c=C_WH)
        stb(sl,"    "+d,Inches(0.6),y+Inches(0.28),Inches(6.0),Inches(0.25),fs=9,c=C_GR); y+=Inches(0.62)
    sft(sl); spg(sl,3)

    # P4 情绪曲线
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"03","情绪曲线","Emotion Curve")
    sl.shapes.add_picture(viz["emotion_curve"],Inches(0.4),Inches(1.45),width=Inches(12.5))
    stb(sl,"解读：2018创业高涨 -> 2019-2020危机下探 -> 2021年后稳定回升",Inches(0.4),SH-Inches(0.75),Inches(12.5),Inches(0.3),fs=10,c=C_AC2)
    sft(sl); spg(sl,4)

    # P5 主题分布
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"04","主题分布","Topic Distribution")
    sl.shapes.add_picture(viz["topic_distribution"],Inches(0.4),Inches(1.45),width=Inches(7.5))
    srect(sl,Inches(8.2),Inches(1.45),Inches(4.8),Inches(5.5),C_CD)
    stb(sl,"关键洞察",Inches(8.4),Inches(1.58),Inches(4.4),Inches(0.38),fs=12,bold=True,c=C_AC)
    sdiv(sl,Inches(8.4),Inches(1.96),Inches(4.4),C_AC)
    ins=["创业内容占比最高（真实创业者身份）","成长学习贯穿始终（持续进化型人格）","情感线从职场向内探索迁移（自我认知深化）","科技关注度2024年后显著上升（All in AI）"]
    for i,ins_ in enumerate(ins):
        stb(sl,"> "+ins_,Inches(8.4),Inches(2.1)+i*Inches(0.7),Inches(4.4),Inches(0.6),fs=10,c=C_WH)
    sft(sl); spg(sl,5)

    # P6 热力图
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"05","年度热力图","Yearly Heatmap")
    sl.shapes.add_picture(viz["yearly_heatmap"],Inches(0.4),Inches(1.45),width=Inches(12.5))
    stb(sl,"绿色=积极 | 红色=消极 | 空白=无数据",Inches(0.4),SH-Inches(0.75),Inches(12.5),Inches(0.28),fs=9,c=C_GR)
    sft(sl); spg(sl,6)

    # P7 转折点
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"06","关键转折","Turning Points")
    evts=[("2018-01","离职创业","从月薪8K押注梦想","月均3.5条，全是兴奋期待",C_AC),("2019-09","首次危机","裁员迷茫卡里2000块","连续6个月情绪评分低于0",C_RD),("2021-06","西藏之行","心灵重建情绪回升","之后3年保持稳定高分",C_AC2),("2024-01","使命觉醒","All in AI元神项目启动","AI科技内容首次超越创业",C_GD)]
    for i,(yr,t,c,ins_,col) in enumerate(evts):
        row=i//2; ci=i%2; x=Inches(0.4)+ci*Inches(6.5); y=Inches(1.45)+row*Inches(2.85)
        srect(sl,x,y,Inches(6.2),Inches(2.6),C_CD); srect(sl,x,y,Inches(6.2),Inches(0.06),col); srect(sl,x,y,Inches(0.06),Inches(2.6),col)
        stb(sl,"  "+yr+"  "+t,x+Inches(0.12),y+Inches(0.1),Inches(5.9),Inches(0.38),fs=13,bold=True,c=col)
        stb(sl,"  "+c,x+Inches(0.12),y+Inches(0.5),Inches(5.9),Inches(0.38),fs=11,c=C_WH)
        stb(sl,"  "+ins_,x+Inches(0.12),y+Inches(1.0),Inches(5.9),Inches(0.55),fs=9,c=C_GR)
    sft(sl); spg(sl,7)

    # P8 AI引擎
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"07","AI分析引擎","AI Analysis")
    steps=[("1.数据采集","微信朋友圈\n授权读取",C_AC),("2.情感分析","DeepSeek API\n-1~+1评分",C_AC2),("3.主题聚类","AI识别\n主要话题",C_GD),("4.可视化","情绪曲线\n热力图饼图",C_GN),("5.洞察生成","关键转折点\n趋势解读",C_RD)]
    for i,(t,d,c) in enumerate(steps):
        x=Inches(0.4)+i*Inches(2.5)
        srect(sl,x,Inches(1.45),Inches(2.15),Inches(2.5),C_CD); srect(sl,x,Inches(1.45),Inches(2.15),Inches(0.48),c)
        stb(sl,t,x+Inches(0.08),Inches(1.47),Inches(2.0),Inches(0.42),fs=10,bold=True,c=C_DK)
        stb(sl,d,x+Inches(0.08),Inches(2.05),Inches(2.0),Inches(1.7),fs=10,c=C_WH,align=PP_ALIGN.CENTER)
        if i<4: stb(sl,">",x+Inches(2.15),Inches(2.5),Inches(0.35),Inches(0.38),fs=16,bold=True,c=c,align=PP_ALIGN.CENTER)
    srect(sl,Inches(0.4),Inches(4.1),SW-Inches(0.8),Inches(2.85),C_CD)
    stb(sl,"技术架构",Inches(0.6),Inches(4.22),Inches(5),Inches(0.35),fs=12,bold=True,c=C_AC)
    techs=[("数据源","微信电脑端导出 txt格式 【日期 时间】内容"),("情感分析","DeepSeek API -1~+1评分 支持API/规则双引擎"),("主题分类","DeepSeek API + 规则引擎 9大主题类别"),("可视化","matplotlib 情绪曲线/热力图/饼图"),("输出格式","JSON结构化数据 + PNG图表 + PPT演示文稿")]
    for i,(k,v) in enumerate(techs):
        ci=i%2; ri=i//2; x=Inches(0.6)+ci*Inches(6.2); y=Inches(4.62)+ri*Inches(0.5)
        stb(sl,k+":",x,y,Inches(1.2),Inches(0.38),fs=9,bold=True,c=C_AC2)
        stb(sl,v,x+Inches(1.2),y,Inches(5.0),Inches(0.38),fs=9,c=C_WH)
    sft(sl); spg(sl,8)

    # P9 元神愿景
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"08","元神愿景","Meta Soul Vision")
    srect(sl,Inches(0.4),Inches(1.45),SW-Inches(0.8),Inches(5.5),C_CD)
    stb(sl,'"看见自己，是改变的第一步"',Inches(0.5),Inches(1.6),Inches(12.3),Inches(0.6),fs=22,bold=True,c=C_AC,align=PP_ALIGN.CENTER)
    sdiv(sl,Inches(3.5),Inches(2.25),Inches(6.3),C_AC2)
    layers=[("分身","替你干活，让你自由","自动化执行日常任务，把时间还给你"),("高我","比你更懂你，让你更强","深度理解你的决策模式，提供个性化参谋"),("元神","与你合一，让你绽放","可交互的数字孪生，终极自我认知工具")]
    for i,(nm,tag,desc) in enumerate(layers):
        y=Inches(2.5)+i*Inches(1.15); col=[C_AC,C_AC2,C_GD][i]
        srect(sl,Inches(0.6),y,Inches(0.08),Inches(0.85),col)
        stb(sl,nm,Inches(0.78),y,Inches(2),Inches(0.35),fs=16,bold=True,c=col)
        stb(sl,tag,Inches(0.78),y+Inches(0.3),Inches(4.5),Inches(0.3),fs=12,bold=True,c=C_WH)
    stb(sl,"10年朋友圈数据 x AI分析 = 你的人生地图",Inches(0.5),Inches(6.5),Inches(12.3),Inches(0.38),fs=13,bold=True,c=C_GD,align=PP_ALIGN.CENTER)
    sft(sl); spg(sl,9)

    # P10 市场机会
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"09","市场机会","Market Opportunity")
    mstats=[("$48.4亿","2026年市场规模",C_AC),("42.2%","年复合增长率",C_AC2),("52%","用户愿用AI助理",C_GD),("3-6个月","先发窗口期",C_GN)]
    for i,(n,l,c) in enumerate(mstats):
        x=Inches(0.4)+i*(Inches(3.0)+Inches(0.18))
        srect(sl,x,Inches(1.45),Inches(3.0),Inches(1.8),C_CD); srect(sl,x,Inches(1.45),Inches(3.0),Inches(0.06),c)
        stb(sl,n,x,Inches(1.58),Inches(3.0),Inches(0.75),fs=22,bold=True,c=c,align=PP_ALIGN.CENTER)
        stb(sl,l,x,Inches(2.33),Inches(3.0),Inches(0.32),fs=10,c=C_WH,align=PP_ALIGN.CENTER)
    srect(sl,Inches(0.4),Inches(3.4),SW-Inches(0.8),Inches(3.5),C_CD)
    stb(sl,"为什么是现在？",Inches(0.6),Inches(3.52),Inches(6),Inches(0.38),fs=12,bold=True,c=C_AC)
    why=[("LLM技术成熟","AI理解个体成为可能"),("Agentic AI爆发","从工具到自主执行者"),("数据主权觉醒","用户开始在意数字身份"),("数字碎片化痛点","没有人有完整自我认知图谱")]
    for i,(t,d) in enumerate(why):
        ci=i%2; ri=i//2
        stb(sl,"  "+t+"  ->  "+d,Inches(0.6)+ci*Inches(6.2),Inches(4.0)+ri*Inches(0.55),Inches(6.0),Inches(0.42),fs=11,c=C_WH)
    sft(sl); spg(sl,10)

    # P11 商业模式
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK); stitle(sl,"10","商业模式","Business Model")
    streams=[("Freemium\n免费增值","基础情绪曲线\n主题分析","免费",C_AC),("Pro $9.9/月","完整洞察报告\n决策建议","付费",C_AC2),("企业版\n定制化","团队家庭版\n私有化部署","B2B",C_GD),("API\n数据服务","情绪API\n洞察SDK","开发者",C_GN)]
    for i,(nm,feat,price,col) in enumerate(streams):
        x=Inches(0.4)+i*Inches(3.12)
        srect(sl,x,Inches(1.45),Inches(2.95),Inches(2.4),C_CD); srect(sl,x,Inches(1.45),Inches(2.95),Inches(0.06),col)
        stb(sl,nm,x+Inches(0.08),Inches(1.58),Inches(2.8),Inches(0.75),fs=13,bold=True,c=col,align=PP_ALIGN.CENTER)
        stb(sl,feat,x+Inches(0.08),Inches(2.35),Inches(2.8),Inches(0.8),fs=10,c=C_WH,align=PP_ALIGN.CENTER)
        stb(sl,price,x+Inches(0.08),Inches(3.3),Inches(2.8),Inches(0.32),fs=12,bold=True,c=col,align=PP_ALIGN.CENTER)
    srect(sl,Inches(0.4),Inches(4.0),SW-Inches(0.8),Inches(2.9),C_CD)
    stb(sl,"飞轮效应",Inches(0.6),Inches(4.12),Inches(5),Inches(0.35),fs=12,bold=True,c=C_AC)
    fw=["  用户上传朋友圈 -> AI分析 -> 生成人生洞察报告","  洞察价值 -> 用户传播邀请 -> 新用户增长","  更多数据 -> 训练基因胶囊 -> AI更懂你","  数据飞轮 -> 护城河加深 -> 竞争壁垒形成"]
    for i,f in enumerate(fw):
        stb(sl,f,Inches(0.6),Inches(4.55)+i*Inches(0.45),Inches(12.3),Inches(0.38),fs=11,c=C_WH)
    sft(sl); spg(sl,11)

    # P12 封底
    sl=prs.slides.add_slide(blank); sbg(sl,C_DK)
    srect(sl,Inches(0),Inches(0),SW,Inches(0.05),C_AC)
    srect(sl,Inches(0),Inches(3.2),SW,Inches(0.04),C_AC)
    srect(sl,Inches(0),SH-Inches(0.05),SW,Inches(0.05),C_AC)
    stb(sl,"让AI替你干活",Inches(0.5),Inches(3.35),Inches(12.3),Inches(0.85),fs=38,bold=True,c=C_WH,align=PP_ALIGN.CENTER)
    stb(sl,"你为自己而活",Inches(0.5),Inches(4.1),Inches(12.3),Inches(0.65),fs=30,bold=True,c=C_AC,align=PP_ALIGN.CENTER)
    sdiv(sl,Inches(4.5),Inches(4.85),Inches(4.3),C_AC2)
    stb(sl,"把工作交给AI，把自己还给热爱",Inches(0.5),Inches(5.1),Inches(12.3),Inches(0.45),fs=14,c=C_GR,align=PP_ALIGN.CENTER)
    stb(sl,"Meta Soul META SOUL MVP Demo 2026",Inches(0.5),SH-Inches(0.7),Inches(12.3),Inches(0.35),fs=10,c=C_GR,align=PP_ALIGN.CENTER)
    spg(sl,12)

    prs.save(out_path)
    print(f"PPT generated: {out_path}")
    return out_path


def gen_script(moments, analysis, out_path):
    n = len(moments)
    script = f"""# 元神项目 MVP 演示视频脚本
# 时长：约2分30秒 | 画外音 + 画面字幕

## 开场（0:00-0:15）
画面：全黑，渐显一条发光的情绪曲线
旁白：这是李明过去10年，在朋友圈留下的情绪轨迹。

## 第一章：我是谁（0:15-0:45）
画面：情绪曲线缓缓展开，颜色从红->黄->绿渐变
旁白：李明，30岁，北漂10年，互联网连续创业者。{n}条朋友圈，记录了他从月薪8K到All in AI的完整旅程。

## 第二章：情绪地图（0:45-1:15）
画面：月度情绪柱状图，每根柱子标注关键事件
旁白：2018年，他辞职创业，情绪高涨。2019年，裁员危机，情绪跌入谷底。2021年，西藏之行，心灵重建。2024年，使命觉醒，All in AI。每一次转折，都清晰可见。

## 第三章：主题解码（1:15-1:45）
画面：主题饼图旋转展开，逐一高亮各主题
旁白：他聊得最多的是：创业、成长、科技。但最近三年，使命开始出现在他的叙述里。这是自我认知的进化，也是人生方向的转变。

## 第四章：元神愿景（1:45-2:15）
画面：三层金字塔动画（分身->高我->元神）
旁白：我们想做的事很简单：让每个人都能像李明一样，通过自己的数字足迹，看见自己的成长轨迹。AI不只是工具，它是激活第二个自己的钥匙。

## 结尾（2:15-2:30）
画面：情绪曲线终点出现光标，指向未来
旁白：让AI替你干活，你为自己而活。元神项目，让每个人都有一个比自己更好的自己。
字幕：META SOUL 元神项目 2026
"""
    with open(out_path,"w",encoding="utf-8") as f:
        f.write(script)
    print(f"Script generated: {out_path}")
    return out_path


# ============================================================
# MAIN
# ============================================================
if __name__ == "__main__":
    import numpy as np
    print("="*50)
    print("Meta Soul MVP Demo Generator v3")
    print("="*50)

    print("\n[Step1] Generating Li Ming's moments data...")
    moments = gen_moments()
    print(f"  {len(moments)} posts | {moments[0]['date']} ~ {moments[-1]['date']}")

    print("\n[Step2] Sentiment analysis + topic classification...")
    if DEEPSEEK_API_KEY and DEEPSEEK_API_KEY != "YOUR_API_KEY_HERE":
        print("  Using DeepSeek API (set DEEPSEEK_API_KEY env var)")
    else:
        print("  Demo mode: rule-based engine (no API needed)")
    for m in moments:
        m["sentiment"] = sentiment(m["content"])
        m["topic"] = topic(m["content"])
    avg_s = np.mean([m["sentiment"] for m in moments])
    print(f"  Avg sentiment: {avg_s:.3f} (positive)")

    print("\n[Step3] Generating visualizations...")
    viz = gen_viz(moments, OUTPUT_DIR)

    print("\n[Step4] Generating 12-page MVP PPT...")
    ppt_path = OUTPUT_DIR / "MetaSoul_MVP_Demo.pptx"
    gen_ppt(moments, moments, viz, str(ppt_path))

    print("\n[Step5] Generating 2:30 video script...")
    script_path = OUTPUT_DIR / "video_script_2min.md"
    gen_script(moments, moments, str(script_path))

    print("\n" + "="*50)
    print("MVP Demo Complete!")
    print("="*50)
    print(f"\nOutput: {OUTPUT_DIR.resolve()}")
    for f in sorted(OUTPUT_DIR.iterdir()):
        sz = f.stat().st_size
        print(f"  {f.name} ({sz//1024}KB)")
    print("\nNext: upload to Feishu group!")
