#!/usr/bin/env /usr/bin/python3
# -*- coding: utf-8 -*-
"""元神AI BP V1.1 PPT生成器 - 深蓝科技风格"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.util import Inches
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image as PILImage

from lxml import etree
import math

# ============================================================
# 配色
# ============================================================
C_BG    = RGBColor(0x0A, 0x19, 0x29)   # 深空蓝背景
C_WHITE = RGBColor(0xF8, 0xFA, 0xFC)   # 星光白
C_GOLD  = RGBColor(0xD4, 0xAF, 0x37)   # 因果金
C_BLUE  = RGBColor(0x21, 0x96, 0xF3)   # 信任蓝
C_GREEN = RGBColor(0x4C, 0xAF, 0x50)   # 生命绿
C_GRAY  = RGBColor(0x94, 0xA3, 0xB8)   # 次要灰
C_CARD  = RGBColor(0x0F, 0x27, 0x47)   # 卡片底色
C_LINE  = RGBColor(0x1E, 0x3A, 0x5F)   # 分割线
C_DIMGOLD = RGBColor(0xB8, 0x96, 0x27) # 深金色

SW = Inches(13.33)
SH = Inches(7.5)

# ============================================================
# 工具函数
# ============================================================
def bg(sl, img_path=None):
    """设置背景：使用预缩放好的1280x720 PNG"""
    if img_path:
        try:
            pic = sl.shapes.add_picture(img_path, Inches(0), Inches(0), width=SW, height=SH)
            # 将图片移到最底层（在文字下方）
            from pptx.oxml.ns import qn
            pic_elem = pic._element
            sp_tree = sl.shapes._spTree
            sp_tree.remove(pic_elem)
            # 找到第一个非nvGrpSpPr/grpSpPr的位置（即文字形状开始的位置）
            insert_pos = 2  # 默认插在 nvGrpSpPr(0) 和 grpSpPr(1) 之后
            for i, child in enumerate(sp_tree):
                tag = child.tag
                if qn('p:sp') in tag or qn('p:pic') in tag:
                    insert_pos = i
                    break
            sp_tree.insert(insert_pos, pic_elem)
        except Exception as e:
            print(f"背景图: {e}, 使用纯色")
            f = sl.background.fill; f.solid(); f.fore_color.rgb = C_BG
    else:
        f = sl.background.fill; f.solid(); f.fore_color.rgb = C_BG

def txt(sl, text, l, t, w, h, fs=14, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False, vAlign=None):
    box = sl.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if vAlign:
        tf.vertical_anchor = vAlign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def rect(sl, l, t, w, h, fill, line=None):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def hline(sl, l, t, w, color=C_GOLD, width=Pt(1.5)):
    ln = sl.shapes.add_shape(1, l, t, w, width)
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln

def card(sl, l, t, w, h, border_color=None):
    r = rect(sl, l, t, w, h, C_CARD, border_color or C_LINE)

def accent_bar(sl, l, t, h, color=C_GOLD, width=Inches(0.06)):
    rect(sl, l, t, width, h, color)

def shadow_box(sl, l, t, w, h, fill, accent=C_GOLD):
    rect(sl, l, t, w, h, fill, C_LINE)
    rect(sl, l, t, w, Inches(0.05), accent)

def page_num(sl, n, total=14):
    txt(sl, f"{n}/{total}", SW - Inches(1.2), SH - Inches(0.38),
        Inches(1), Inches(0.25), fs=12, color=C_GRAY, align=PP_ALIGN.RIGHT)

def footer(sl):
    txt(sl, "元神AI · 个人因果模型 · 2026",
        Inches(0.4), SH - Inches(0.38), Inches(5), Inches(0.25),
        fs=12, color=C_GRAY)

def section_tag(sl, text, l, t):
    """小标签"""
    rect(sl, l, t, Inches(len(text) * 0.11 + 0.3), Inches(0.32), C_GOLD)
    txt(sl, text, l + Inches(0.1), t + Inches(0.03), Inches(4), Inches(0.28),
        fs=9, bold=True, color=C_BG)

# ============================================================
# 幻灯片生成
# ============================================================
def make_cover(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg1_1280.png")

    # 顶部金色细线
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))
    # 底部金色细线
    hline(sl, Inches(0), SH - Inches(0.08), SW, C_GOLD, Pt(2))
    # 中间装饰分隔线
    hline(sl, Inches(3), Inches(3.5), Inches(7.33), C_GOLD, Pt(1.5))
    # 简洁星光装饰
    for i in range(6):
        dot = sl.shapes.add_shape(9, Inches(0.5 + i * 2.2), Inches(0.4),
                                  Inches(0.05), Inches(0.05))
        dot.fill.solid(); dot.fill.fore_color.rgb = C_GOLD
        dot.line.fill.background()
    # 右上角装饰
    for i in range(5):
        dot = sl.shapes.add_shape(9, SW - Inches(0.8), Inches(0.3 + i * 0.2),
                                  Inches(0.04), Inches(0.04))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_GOLD
        dot.line.fill.background()

    page_num(sl, 1)


def make_exec_summary(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))

    # 标签
    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), C_LINE, Pt(0.5))

    # 核心句（居中大字）
    txt(sl, "大模型知道世界，但不知道你的世界",
        Inches(0.5), Inches(0.9), Inches(12.33), Inches(0.7),
        fs=44, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 三个金色数字（建议修改项4）
    big_nums = [("500万", "融资目标"), ("4500万", "投前估值"), ("10万（首年目标）", "深度用户")]
    nw = Inches(3.8)
    for i, (n, l) in enumerate(big_nums):
        nx = Inches(0.5) + i * (nw + Inches(0.25))
        txt(sl, n, nx, Inches(1.75), nw, Inches(0.8),
            fs=32, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        txt(sl, l, nx, Inches(2.62), nw, Inches(0.35),
            fs=13, color=C_WHITE, align=PP_ALIGN.CENTER)

    hline(sl, Inches(0.5), Inches(3.05), Inches(12.33), C_LINE, Pt(0.5))

    # 正文摘要
    card(sl, Inches(0.5), Inches(3.15), Inches(12.33), Inches(2.8))
    para = (
        "我们正处在一个AI爆发却人类愈发孤独的时代：信息爆炸带来决策过载，"
        "算法推荐制造认知茧房，而最懂你的那个人——你自己——却从未被真正\"看见\"。\n\n"
        "元神AI切入一个被所有人忽略的赛道：**个人因果模型**。"
        "不是给AI装一个你的记忆，而是让AI理解你为什么会做A而不做B。"
        "我们将每个人的精神世界建模为一张因果网络，让AI在高我层给你参谋，"
        "在分身层替你执行，在元神层与你合一。\n\n"
        "三年内，让个人因果模型成为数字时代的基础设施。"
    )
    txt(sl, para, Inches(0.7), Inches(3.3), Inches(11.93), Inches(2.5),
        fs=14, color=C_WHITE)

    footer(sl)
    page_num(sl, 2)


def make_pain_points(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))

    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), C_LINE, Pt(0.5))

    txt(sl, "市场之痛：双重困境",
        Inches(0.5), Inches(0.85), Inches(12.33), Inches(0.6),
        fs=44, bold=True, color=C_GOLD)

    # 左栏：决策过载
    card(sl, Inches(0.5), Inches(1.65), Inches(5.9), Inches(4.7), C_BLUE)
    rect(sl, Inches(0.5), Inches(1.65), Inches(5.9), Inches(0.05), C_BLUE)
    txt(sl, "困境一", Inches(0.7), Inches(1.78), Inches(2), Inches(0.35),
        fs=10, bold=True, color=C_BLUE)
    txt(sl, "决策过载",
        Inches(0.7), Inches(2.08), Inches(5.5), Inches(0.55),
        fs=22, bold=True, color=C_WHITE)
    txt(sl, "现代人每天面临数百次显性决策（认知心理学共识），"
        "深夜复盘时却发现：当年那个\"正确的选择\"，至今没被真正理解过。"
        "你在每个决策节点留下了数据碎片，但没有人把它们连成你的因果逻辑。",
        Inches(0.7), Inches(2.72), Inches(5.5), Inches(1.8),
        fs=13, color=C_WHITE)
    # 图标示意：数据碎片
    for j in range(5):
        rx = Inches(0.9 + j * 1.05)
        ry = Inches(4.7)
        rw = Inches(0.85)
        rh = Inches(0.55)
        rect(sl, rx, ry, rw, rh, RGBColor(0x1a, 0x35 + j * 8, 0x5c + j * 5))
        txt(sl, f"决策{j+1}", rx, ry + Inches(0.12), rw, Inches(0.3),
            fs=8, color=C_WHITE, align=PP_ALIGN.CENTER)
    hline(sl, Inches(1.2), Inches(5.35), Inches(4.6), C_BLUE, Pt(1))
    txt(sl, "碎片 → 无人连接",
        Inches(0.7), Inches(5.5), Inches(5.5), Inches(0.35),
        fs=11, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    # 右栏：灵魂孤独
    card(sl, Inches(6.9), Inches(1.65), Inches(5.9), Inches(4.7), C_GOLD)
    rect(sl, Inches(6.9), Inches(1.65), Inches(5.9), Inches(0.05), C_GOLD)
    txt(sl, "困境二", Inches(7.1), Inches(1.78), Inches(2), Inches(0.35),
        fs=10, bold=True, color=C_GOLD)
    txt(sl, "灵魂孤独",
        Inches(7.1), Inches(2.08), Inches(5.5), Inches(0.55),
        fs=22, bold=True, color=C_WHITE)
    txt(sl, "不是没有倾诉对象，而是没有\"比你更懂你\"的知己。"
        "家人爱你但不懂你的事业焦虑，朋友陪你但看不见你的内在轨迹，"
        "而现有AI只是工具——问什么答什么，不主动，不懂你。",
        Inches(7.1), Inches(2.72), Inches(5.5), Inches(1.8),
        fs=13, color=C_WHITE)
    # 三个"AI"方块代表现有工具
    for j in range(3):
        rect(sl, Inches(7.3 + j * 1.7), Inches(4.7), Inches(1.4), Inches(0.55),
             RGBColor(0x2a, 0x2a, 0x3a))
        txt(sl, "AI工具", Inches(7.3 + j * 1.7), Inches(4.8),
            Inches(1.4), Inches(0.35),
            fs=9, color=C_GRAY, align=PP_ALIGN.CENTER)
    hline(sl, Inches(7.1), Inches(5.35), Inches(5.5), C_GOLD, Pt(1))
    txt(sl, "问什么答什么 ≠ 真正懂你",
        Inches(7.1), Inches(5.5), Inches(5.5), Inches(0.35),
        fs=11, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 底部总结
    hline(sl, Inches(0.5), Inches(6.5), Inches(12.33), C_LINE, Pt(0.5))
    txt(sl, "现有AI的根本局限：它们是通用秘书，不是你的高参。",
        Inches(0.5), Inches(6.58), Inches(12.33), Inches(0.45),
        fs=13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    footer(sl)
    page_num(sl, 3)


def make_orange_sea(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))

    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), C_LINE, Pt(0.5))

    # 三个圆（Venn图）
    # 红海圆
    cx1 = Inches(3.2); cy = Inches(3.2); r = Inches(1.5)
    c1 = sl.shapes.add_shape(9, cx1 - r, cy - r, r * 2, r * 2)
    c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x3F, 0x1E, 0x3F)
    c1.line.color.rgb = RGBColor(0xEF, 0x44, 0x44)
    c1.line.width = Pt(2)

    # 蓝海圆
    cx2 = Inches(5.2)
    c2 = sl.shapes.add_shape(9, cx2 - r, cy - r, r * 2, r * 2)
    c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x0A, 0x19, 0x40)
    c2.line.color.rgb = RGBColor(0x44, 0x88, 0xEF)
    c2.line.width = Pt(2)

    # 橙海圆（中间重叠）
    cx3 = Inches(4.2); cy3 = Inches(3.5); r3 = Inches(1.2)
    c3 = sl.shapes.add_shape(9, cx3 - r3, cy3 - r3, r3 * 2, r3 * 2)
    c3.fill.solid(); c3.fill.fore_color.rgb = RGBColor(0x2A, 0x18, 0x05)
    c3.line.color.rgb = C_GOLD
    c3.line.width = Pt(2.5)

    # 圆心标签
    txt(sl, "红海", Inches(2.5), Inches(2.8), Inches(1.2), Inches(0.4),
        fs=13, bold=True, color=RGBColor(0xEF, 0x44, 0x44), align=PP_ALIGN.CENTER)
    txt(sl, "更聪明的工具", Inches(2.2), Inches(3.15), Inches(1.8), Inches(0.3),
        fs=9, color=RGBColor(0xCC, 0x88, 0x88), align=PP_ALIGN.CENTER)

    txt(sl, "蓝海", Inches(5.3), Inches(2.8), Inches(1.2), Inches(0.4),
        fs=13, bold=True, color=RGBColor(0x44, 0x88, 0xEF), align=PP_ALIGN.CENTER)
    txt(sl, "无人验证的需求", Inches(5.1), Inches(3.15), Inches(1.6), Inches(0.3),
        fs=9, color=RGBColor(0x88, 0xAA, 0xDD), align=PP_ALIGN.CENTER)

    txt(sl, "橙海 ✓", Inches(3.4), Inches(3.3), Inches(1.6), Inches(0.45),
        fs=16, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
    txt(sl, "更懂你的关系", Inches(3.2), Inches(3.72), Inches(2.0), Inches(0.3),
        fs=11, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 右侧大字
    txt(sl, "我们不在\n红海卷，\n也不在\n蓝海赌",
        Inches(7.5), Inches(1.2), Inches(5.5), Inches(2.5),
        fs=30, bold=True, color=C_WHITE)
    hline(sl, Inches(7.5), Inches(3.75), Inches(5.33), C_GOLD, Pt(2))
    txt(sl, "我们在橙海",
        Inches(7.5), Inches(3.9), Inches(5.5), Inches(0.65),
        fs=48, bold=True, color=C_GOLD)

    # 底部说明
    card(sl, Inches(0.5), Inches(4.9), Inches(12.33), Inches(1.3))
    txt(sl, "现有玩家都在做\"更聪明的工具\"，而我们在做\"更懂你的关系\"。"
        "一个被所有人看见却无人真正理解的市场——这才是元神AI的起点。",
        Inches(0.7), Inches(5.0), Inches(11.93), Inches(1.0),
        fs=13, color=C_WHITE, align=PP_ALIGN.CENTER)

    footer(sl)
    page_num(sl, 4)


def make_vision_overview(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))

    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), C_LINE, Pt(0.5))

    txt(sl, "元神之道：三层次价值",
        Inches(0.5), Inches(0.85), Inches(12.33), Inches(0.6),
        fs=44, bold=True, color=C_GOLD)

    # 三列卡片
    layers = [
        ("分身", "替你干活，让你自由", [
            "处理邮件、日程、信息过滤",
            "把机械性事务接管",
            "让你腾出时间给热爱的事"
        ], C_BLUE, "01"),
        ("高我", "比你更懂你，让你更强", [
            "深度理解你的决策逻辑",
            "价值观排序、情绪轨迹",
            "关键决策前精准参谋"
        ], C_GOLD, "02"),
        ("元神", "与你合一，让你绽放", [
            "构建个人因果模型",
            "可交互、可预测",
            "可持续进化"
        ], C_GREEN, "03"),
    ]

    cw = Inches(3.9)
    for i, (name, tagline, points, col, num) in enumerate(layers):
        cx = Inches(0.5) + i * (cw + Inches(0.25))
        card(sl, cx, Inches(1.65), cw, Inches(4.0), col)
        rect(sl, cx, Inches(1.65), cw, Inches(0.06), col)
        # 编号
        txt(sl, num, cx + Inches(0.15), Inches(1.8), Inches(0.5), Inches(0.5),
            fs=28, bold=True, color=col)
        # 名称
        txt(sl, name, cx + Inches(0.15), Inches(2.35), cw - Inches(0.3), Inches(0.55),
            fs=22, bold=True, color=C_WHITE)
        # 副标题
        txt(sl, tagline, cx + Inches(0.15), Inches(2.9), cw - Inches(0.3), Inches(0.4),
            fs=11, italic=True, color=col)
        hline(sl, cx + Inches(0.15), Inches(3.35), cw - Inches(0.3), C_LINE, Pt(0.5))
        # 要点
        for j, pt in enumerate(points):
            txt(sl, f"› {pt}", cx + Inches(0.15), Inches(3.5) + j * Inches(0.5),
                cw - Inches(0.3), Inches(0.45), fs=12, color=C_WHITE)

    # 底部金句
    hline(sl, Inches(0.5), Inches(5.85), Inches(12.33), C_LINE, Pt(0.5))
    card(sl, Inches(3.5), Inches(5.95), Inches(6.33), Inches(0.9), C_GOLD)
    txt(sl, "工具可以被替代，关系无法被删除。",
        Inches(3.5), Inches(6.05), Inches(6.33), Inches(0.7),
        fs=16, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

    footer(sl)
    page_num(sl, 5)


def _make_layer_slide(prs, num, name, tagline, story, key_point, col, col_dark):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, col, Pt(2))

    # 左侧70% 内容
    section_tag(sl, f"LAYER 0{num}", Inches(0.5), Inches(0.3))
    hline(sl, Inches(0.5), Inches(0.72), Inches(8.5), C_LINE, Pt(0.5))

    txt(sl, name, Inches(0.5), Inches(0.85), Inches(8.5), Inches(0.75),
        fs=44, bold=True, color=C_GOLD)
    txt(sl, tagline, Inches(0.5), Inches(1.6), Inches(8.5), Inches(0.45),
        fs=16, italic=True, color=col)

    hline(sl, Inches(0.5), Inches(2.1), Inches(8.5), col, Pt(1))

    # 用户故事框
    card(sl, Inches(0.5), Inches(2.3), Inches(8.5), Inches(2.8), col)
    txt(sl, "用户故事", Inches(0.7), Inches(2.42), Inches(2), Inches(0.32),
        fs=10, bold=True, color=col)
    txt(sl, story, Inches(0.7), Inches(2.8), Inches(8.1), Inches(2.1),
        fs=13, color=C_WHITE)

    # 底部关键点
    hline(sl, Inches(0.5), Inches(5.25), Inches(8.5), C_LINE, Pt(0.5))
    txt(sl, "核心价值：", Inches(0.5), Inches(5.35), Inches(2), Inches(0.4),
        fs=12, bold=True, color=C_GRAY)
    txt(sl, key_point, Inches(2.4), Inches(5.35), Inches(6.5), Inches(0.5),
        fs=14, bold=True, color=col)

    # 右侧视觉区
    card(sl, Inches(9.4), Inches(0.5), Inches(3.6), Inches(6.3))

    # 抽象图形
    if num == 1:
        # 分身：齿轮/循环箭头
        for k in range(4):
            ang = k * 90
            ax = Inches(10.5) + Inches(math.cos(math.radians(ang)) * 1.0)
            ay = Inches(3.5) + Inches(math.sin(math.radians(ang)) * 0.8)
            dot = sl.shapes.add_shape(9, ax, ay, Inches(0.5), Inches(0.5))
            dot.fill.solid(); dot.fill.fore_color.rgb = col
            dot.line.fill.background()
        dot_c = sl.shapes.add_shape(9, Inches(10.5), Inches(3.5), Inches(0.6), Inches(0.6))
        dot_c.fill.solid(); dot_c.fill.fore_color.rgb = col_dark
        dot_c.line.color.rgb = col; dot_c.line.width = Pt(2)
        txt(sl, "分身", Inches(9.6), Inches(1.2), Inches(3.2), Inches(0.5),
            fs=20, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, "替你干活\n让你自由", Inches(9.6), Inches(5.0), Inches(3.2), Inches(1.0),
            fs=14, color=C_GRAY, align=PP_ALIGN.CENTER)

    elif num == 2:
        # 高我：双向箭头/镜子
        for k in range(3):
            bx = Inches(10.0) + k * Inches(0.6)
            rect(sl, bx, Inches(2.5) + k * Inches(0.6), Inches(0.3), Inches(1.2),
                 RGBColor(0x0a + k * 8, 0x25 + k * 6, 0x45))
        arrow_h = sl.shapes.add_shape(1, Inches(9.9), Inches(4.2), Inches(3.2), Inches(0.06))
        arrow_h.fill.solid(); arrow_h.fill.fore_color.rgb = col; arrow_h.line.fill.background()
        txt(sl, "高我", Inches(9.6), Inches(1.2), Inches(3.2), Inches(0.5),
            fs=20, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, "比你更懂你\n让你更强", Inches(9.6), Inches(5.0), Inches(3.2), Inches(1.0),
            fs=14, color=C_GRAY, align=PP_ALIGN.CENTER)

    else:
        # 元神：三层同心圆
        for k in range(3):
            r_size = Inches(1.6 - k * 0.4)
            cr = sl.shapes.add_shape(9,
                Inches(10.9) - r_size/2, Inches(3.5) - r_size/2,
                r_size, r_size)
            cr.fill.solid()
            cr.fill.fore_color.rgb = RGBColor(0x0A + k*5, 0x19 + k*5, 0x29 + k*5)
            cr.line.color.rgb = [col, C_BLUE, C_GRAY][k]
            cr.line.width = Pt([2, 1.5, 1][k])
        txt(sl, "元神", Inches(9.6), Inches(1.2), Inches(3.2), Inches(0.5),
            fs=20, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, "与你合一\n让你绽放", Inches(9.6), Inches(5.0), Inches(3.2), Inches(1.0),
            fs=14, color=C_GRAY, align=PP_ALIGN.CENTER)

    footer(sl)
    page_num(sl, 5 + num)


def make_tech(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))

    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), C_LINE, Pt(0.5))

    txt(sl, "技术护城河：四层架构",
        Inches(0.5), Inches(0.85), Inches(12.33), Inches(0.6),
        fs=44, bold=True, color=C_GOLD)

    layers = [
        ("01", "意图识别引擎", "通义千问", "从对话、行为、时间序列中抽取意图向量，"
         "区分嘴上说的和真正做的。", C_BLUE),
        ("02", "因果推理引擎", "LingBot-VA", "不是相关分析，是因果建模。"
         "理解深层因果链。", C_GOLD),
        ("03", "个人世界模型", "LeWorldModel / Emu3", "将因果网络建模为"
         "可推理的嵌入空间，实时更新，支持模拟预测。", RGBColor(0x8B, 0x5C, 0xF6)),
        ("04", "主动对话系统", "主动发起 × 关系深化", "不是等用户提问，而是AI主动发起。"
         "每周内在状态报告，每逢决策主动参谋。", C_GREEN),
    ]

    bh = Inches(1.1)
    for i, (num, name, tech, desc, col) in enumerate(layers):
        by = Inches(1.65) + i * (bh + Inches(0.1))
        # 底层卡片
        card(sl, Inches(0.5), by, Inches(12.33), bh, col)
        rect(sl, Inches(0.5), by, Inches(0.08), bh, col)
        # 编号
        txt(sl, num, Inches(0.65), by + Inches(0.1), Inches(0.5), Inches(0.45),
            fs=20, bold=True, color=col)
        # 名称
        txt(sl, name, Inches(1.2), by + Inches(0.1), Inches(3.5), Inches(0.45),
            fs=16, bold=True, color=C_WHITE)
        # 技术标签
        rect(sl, Inches(1.2), by + Inches(0.58), Inches(len(tech) * 0.1 + 0.3), Inches(0.3), col)
        txt(sl, tech, Inches(1.28), by + Inches(0.6), Inches(4), Inches(0.28),
            fs=9, bold=True, color=C_BG)
        # 描述
        txt(sl, desc, Inches(5.5), by + Inches(0.15), Inches(7.1), Inches(0.85),
            fs=12, color=C_WHITE)
        # 连接箭头
        if i < 3:
            txt(sl, "▼", Inches(6.5), by + bh, Inches(0.4), Inches(0.2),
                fs=14, color=col, align=PP_ALIGN.CENTER)

    footer(sl)
    page_num(sl, 9)


def make_flywheel(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))

    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), C_LINE, Pt(0.5))

    txt(sl, "数据飞轮效应",
        Inches(0.5), Inches(0.85), Inches(12.33), Inches(0.6),
        fs=44, bold=True, color=C_GOLD)

    # 环形飞轮（五节点）
    nodes = [
        ("更多用户", C_BLUE),
        ("更准模型", C_GOLD),
        ("更高价值", C_GREEN),
        ("更强粘性", RGBColor(0x8B, 0x5C, 0xF6)),
        ("更多数据", RGBColor(0xF4, 0x43, 0xA6)),
    ]
    cx = Inches(6.66); cy = Inches(4.0); R = Inches(2.2)
    angles = [90, 18, -54, -126, -198]

    # 画环形箭头（用虚线圆代替箭头，省略箭头复杂性）
    for i, (angle, (label, col)) in enumerate(zip(angles, nodes)):
        rad = math.radians(angle)
        nx = cx + R * math.cos(rad)
        ny = cy - R * math.sin(rad)
        # 节点圆
        node = sl.shapes.add_shape(9, nx - Inches(0.65), ny - Inches(0.45),
                                    Inches(1.3), Inches(0.9))
        node.fill.solid(); node.fill.fore_color.rgb = col
        node.line.fill.background()
        txt(sl, label, nx - Inches(0.65), ny - Inches(0.28),
            Inches(1.3), Inches(0.4),
            fs=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        # 连接弧（用细线段近似）
        next_rad = math.radians(angles[(i+1) % 5])
        lx1 = nx + Inches(0.5) * math.cos(rad)
        ly1 = ny - Inches(0.5) * math.sin(rad)
        lx2 = cx + (R - Inches(0.1)) * math.cos(next_rad)
        ly2 = cy - (R - Inches(0.1)) * math.sin(next_rad)
        ln = sl.shapes.add_shape(1, min(lx1, lx2), min(ly1, ly2),
                                  abs(lx2 - lx1) + Inches(0.02), abs(ly2 - ly1) + Inches(0.02))
        ln.fill.background(); ln.line.color.rgb = C_GRAY; ln.line.width = Pt(1)
        # 箭头符号
        ax = (nx + lx2) / 2; ay = (ny + ly2) / 2
        txt(sl, "→", ax, ay - Inches(0.15), Inches(0.3), Inches(0.3),
            fs=14, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 中心
    cen = sl.shapes.add_shape(9, cx - Inches(0.55), cy - Inches(0.4),
                               Inches(1.1), Inches(0.8))
    cen.fill.solid(); cen.fill.fore_color.rgb = C_CARD
    cen.line.color.rgb = C_GOLD; cen.line.width = Pt(2)
    txt(sl, "元神AI", cx - Inches(0.55), cy - Inches(0.2),
        Inches(1.1), Inches(0.4),
        fs=14, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    # 右侧说明
    card(sl, Inches(10.0), Inches(1.65), Inches(2.9), Inches(5.3))
    txt(sl, "飞轮逻辑", Inches(10.15), Inches(1.8), Inches(2.6), Inches(0.4),
        fs=13, bold=True, color=C_GOLD)
    hline(sl, Inches(10.0), Inches(2.2), Inches(2.9), C_LINE, Pt(0.5))
    points = [
        ("用户", "上传朋友圈\n/聊天数据"),
        ("模型", "训练因果\n推理网络"),
        ("价值", "生成精准\n洞察报告"),
        ("粘性", "用户依赖\n主动对话"),
        ("数据", "更多数据\n反哺模型"),
    ]
    for j, (k, v) in enumerate(points):
        txt(sl, f"{j+1}. {k}", Inches(10.15), Inches(2.35) + j * Inches(0.85),
            Inches(2.6), Inches(0.28),
            fs=11, bold=True, color=C_WHITE)
        txt(sl, v, Inches(10.15), Inches(2.6) + j * Inches(0.85),
            Inches(2.6), Inches(0.55), fs=9, color=C_GRAY)

    # 底部金句
    hline(sl, Inches(0.5), Inches(6.5), Inches(12.33), C_LINE, Pt(0.5))
    txt(sl, "护城河随时间加深 · 先发优势一旦建立，难以被追赶",
        Inches(0.5), Inches(6.58), Inches(12.33), Inches(0.45),
        fs=14, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    footer(sl)
    page_num(sl, 10)


def make_biz_model(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))

    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), C_LINE, Pt(0.5))

    txt(sl, "商业模式：三位一体",
        Inches(0.5), Inches(0.85), Inches(12.33), Inches(0.6),
        fs=44, bold=True, color=C_GOLD)

    plans = [
        ("FREE", "免费层", "引流入口", [
            "情绪日报",
            "每周洞察摘要",
            "基础问答",
        ], C_GRAY, "永久免费"),
        ("PRO", "Pro 订阅", "398元/年（39元/月）", [
            "完整因果报告",
            "决策参谋",
            "主动对话",
            "年度人生复盘",
        ], C_BLUE, "249元/年"),
        ("SHARE", "按果分成", "核心创新 ★", [
            "重大决策归因",
            "5-15% 价值分成",
            "用户自主设定比例",
            "可验证性基于因果归因引擎",
            "——记录决策前后因果链变化，生成可审计报告，分成由用户自主确认",
        ], C_GOLD, "因果价值 × 成果"),
    ]

    pw = Inches(4.0)
    for i, (tag, name, price, feats, col, sub) in enumerate(plans):
        px = Inches(0.5) + i * (pw + Inches(0.25))
        card(sl, px, Inches(1.65), pw, Inches(4.8), col)
        rect(sl, px, Inches(1.65), pw, Inches(0.06), col)
        # 标签
        rect(sl, px + Inches(0.15), Inches(1.78), Inches(1.0), Inches(0.32), col)
        txt(sl, tag, px + Inches(0.15), Inches(1.8), Inches(1.0), Inches(0.28),
            fs=10, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        # 名称
        txt(sl, name, px + Inches(0.15), Inches(2.2), pw - Inches(0.3), Inches(0.55),
            fs=18, bold=True, color=C_WHITE)
        # 价格
        txt(sl, price, px + Inches(0.15), Inches(2.75), pw - Inches(0.3), Inches(0.5),
            fs=14, bold=True, color=col)
        # 分隔线
        hline(sl, px + Inches(0.15), Inches(3.28), pw - Inches(0.3), C_LINE, Pt(0.5))
        # 功能列表
        for j, f in enumerate(feats):
            txt(sl, f"✓  {f}", px + Inches(0.15), Inches(3.42) + j * Inches(0.4),
                pw - Inches(0.3), Inches(0.38), fs=12, color=C_WHITE)
        # 底部说明
        hline(sl, px + Inches(0.15), Inches(5.65), pw - Inches(0.3), C_LINE, Pt(0.5))
        txt(sl, sub, px + Inches(0.15), Inches(5.75), pw - Inches(0.3), Inches(0.45),
            fs=10, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)

    # 底部飞轮提示
    hline(sl, Inches(0.5), Inches(6.6), Inches(12.33), C_LINE, Pt(0.5))
    txt(sl, "分成比例由用户自主设定，我们只收取用户认可的\"因果价值\"",
        Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.4),
        fs=12, italic=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    footer(sl)
    page_num(sl, 11)


def make_founder(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))

    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), C_LINE, Pt(0.5))

    txt(sl, "创始人与使命：十年藏心",
        Inches(0.5), Inches(0.85), Inches(12.33), Inches(0.6),
        fs=44, bold=True, color=C_GOLD)

    # 左侧：排比困境（扉页格式）
    card(sl, Inches(0.5), Inches(1.65), Inches(6.2), Inches(5.2))
    # 扉页格式首句
    txt(sl, "我曾深夜在阳台追问：",
        Inches(0.7), Inches(1.78), Inches(5.8), Inches(0.38),
        fs=11, color=C_GRAY)
    txt(sl, "为什么我读了那么多书，还是过不好这一生？",
        Inches(0.7), Inches(2.2), Inches(5.8), Inches(0.55),
        fs=13, italic=True, color=C_WHITE)
    qs = [
        "为什么我做了那么多正确的事，\n还是不快乐？",
        "为什么我明明有很多朋友，\n却依然感到孤独？",
    ]
    for j, q in enumerate(qs):
        txt(sl, '"', Inches(0.6), Inches(2.88) + j * Inches(1.05),
            Inches(0.4), Inches(0.5),
            fs=44, bold=True, color=C_GOLD)
        txt(sl, q, Inches(1.0), Inches(2.88) + j * Inches(1.05),
            Inches(5.5), Inches(0.85),
            fs=13, italic=True, color=C_WHITE)
    hline(sl, Inches(0.7), Inches(5.0), Inches(5.8), C_GOLD, Pt(1))
    txt(sl, "后来我才明白——因为我一直向外求。",
        Inches(0.7), Inches(5.1), Inches(5.8), Inches(0.55),
        fs=13, bold=True, color=C_GOLD)

    # 右侧：时间线
    card(sl, Inches(7.0), Inches(1.65), Inches(5.9), Inches(5.2))
    txt(sl, "十年轨迹", Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.4),
        fs=13, bold=True, color=C_GOLD)
    hline(sl, Inches(7.2), Inches(2.2), Inches(5.5), C_LINE, Pt(0.5))

    milestones = [
        ("2015", "出发", "北京出租屋读《原则》，\n月薪8K辞职创业", C_BLUE),
        ("2018", "创办", "创业者夜校", C_GRAY),
        ("2019", "低谷", "公司失败，账上2000块，\n深夜朋友圈收到温暖回复", RGBColor(0xEF, 0x44, 0x44)),
        ("2020", "创办", "大思维智库", C_GRAY),
        ("2021", "重建", "西藏布达拉宫前哭2小时，\n放下对\"成功\"的执念", C_GREEN),
        ("2024", "觉醒", "All in AI，元神项目启动，\n答案清晰：帮人更懂自己", C_GOLD),
    ]
    for j, (yr, t, d, col) in enumerate(milestones):
        yy = Inches(2.35) + j * Inches(0.7)
        # 时间线竖线
        rect(sl, Inches(7.35), yy, Inches(0.04), Inches(0.75), col)
        # 年份
        txt(sl, yr, Inches(7.5), yy, Inches(1.0), Inches(0.32),
            fs=12, bold=True, color=col)
        # 标题
        txt(sl, t, Inches(8.55), yy, Inches(4.2), Inches(0.3),
            fs=11, bold=True, color=C_WHITE)
        # 描述
        txt(sl, d, Inches(8.55), yy + Inches(0.28), Inches(4.2), Inches(0.48),
            fs=9, color=C_GRAY)

    footer(sl)
    page_num(sl, 12)


def make_funding(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg3_1280.png")
    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(2))

    hline(sl, Inches(0.5), Inches(0.72), Inches(12.33), C_LINE, Pt(0.5))

    txt(sl, "融资与用途：种子轮 500万",
        Inches(0.5), Inches(0.85), Inches(12.33), Inches(0.6),
        fs=44, bold=True, color=C_GOLD)

    # 三个核心数据
    stats = [("4500万", "投前估值", C_GOLD), ("500万", "融资目标", C_BLUE), ("10%", "股权释放", C_GREEN)]
    for i, (n, l, col) in enumerate(stats):
        sx = Inches(0.5) + i * Inches(4.25)
        card(sl, sx, Inches(1.6), Inches(3.9), Inches(1.3), col)
        rect(sl, sx, Inches(1.6), Inches(3.9), Inches(0.05), col)
        txt(sl, n, sx, Inches(1.72), Inches(3.9), Inches(0.65),
            fs=28, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)
        txt(sl, l, sx, Inches(2.42), Inches(3.9), Inches(0.35),
            fs=12, color=C_WHITE, align=PP_ALIGN.CENTER)

    # 饼图文字描述
    card(sl, Inches(0.5), Inches(3.1), Inches(5.8), Inches(3.6))
    txt(sl, "资金分配", Inches(0.7), Inches(3.22), Inches(5.4), Inches(0.4),
        fs=14, bold=True, color=C_GOLD)
    hline(sl, Inches(0.5), Inches(3.65), Inches(5.8), C_LINE, Pt(0.5))

    allocations = [
        ("40%", "产品研发\n200万 · 因果引擎 + 主动对话系统", C_GOLD, Inches(0.52)),
        ("30%", "团队薪酬\n150万 · 核心技术人员与产品经理", C_BLUE, Inches(0.52 + 1.3)),
        ("20%", "市场获客\n100万 · 种子用户 + KOL测评 + 社群", C_GREEN, Inches(0.52 + 2.6)),
        ("10%", "运营储备\n50万 · 云服务 + API + 法务合规", C_GRAY, Inches(0.52 + 3.9)),
    ]
    for pct, desc, col, py in allocations:
        rect(sl, Inches(0.7), py, Inches(0.7), Inches(0.5), col)
        txt(sl, pct, Inches(0.7), py + Inches(0.05), Inches(0.7), Inches(0.4),
            fs=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        txt(sl, desc, Inches(1.5), py, Inches(4.6), Inches(0.65),
            fs=11, color=C_WHITE)

    # 嵌入已生成的饼图
    sl.shapes.add_picture("/Users/w/.openclaw/workspace/funding_pie.png",
                          Inches(6.6), Inches(3.1), width=Inches(6.4))

    # 里程碑
    hline(sl, Inches(0.5), Inches(6.75), Inches(12.33), C_LINE, Pt(0.5))
    milestones_bar = [
        ("3个月", "MVP上线", C_BLUE),
        ("12个月", "1万付费用户", C_GOLD),
        ("24个月", "盈亏平衡", C_GREEN),
    ]
    for i, (t, m, col) in enumerate(milestones_bar):
        mx = Inches(0.5) + i * Inches(4.25)
        rect(sl, mx, Inches(6.85), Inches(3.9), Inches(0.05), col)
        txt(sl, t, mx, Inches(6.95), Inches(1.5), Inches(0.35),
            fs=12, bold=True, color=col)
        txt(sl, m, mx + Inches(1.5), Inches(6.95), Inches(2.4), Inches(0.35),
            fs=12, color=C_WHITE)

    footer(sl)
    page_num(sl, 13)


def make_ending(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, "/Users/w/.openclaw/workspace/bg1_1280.png")

    hline(sl, Inches(0), Inches(0), SW, C_GOLD, Pt(3))
    hline(sl, Inches(0), Inches(3.1), SW, C_GOLD, Pt(1.5))
    hline(sl, Inches(0), SH - Inches(0.1), SW, C_GOLD, Pt(3))

    # 简洁星光装饰
    for i in range(6):
        dot = sl.shapes.add_shape(9, Inches(0.5 + i * 2.2), Inches(0.3),
                                  Inches(0.06), Inches(0.06))
        dot.fill.solid(); dot.fill.fore_color.rgb = C_GOLD
        dot.line.fill.background()

    txt(sl, "让AI替你干活",
        Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.9),
        fs=44, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

    txt(sl, "你为自己而活",
        Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.7),
        fs=36, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    hline(sl, Inches(3), Inches(4.88), Inches(7.33), C_GOLD, Pt(1))

    txt(sl, "把生活交给AI，把人生还给热爱",
        Inches(0.5), Inches(5.05), Inches(12.33), Inches(0.5),
        fs=16, color=C_GRAY, align=PP_ALIGN.CENTER)

    txt(sl, "元神AI · 让每个人都有一个比自己更好的自己",
        Inches(0.5), SH - Inches(0.95), Inches(9), Inches(0.4),
        fs=11, color=C_GOLD, align=PP_ALIGN.CENTER)

    txt(sl, "吴治广  创始人 · 元神AI",
        SW - Inches(4.5), SH - Inches(0.95), Inches(4), Inches(0.35),
        fs=9, color=C_GRAY, align=PP_ALIGN.RIGHT)

    page_num(sl, 14)


# ============================================================
# 主程序
# ============================================================
if __name__ == "__main__":
    print("开始生成 BP V1.1 PPT...")

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    make_cover(prs)
    print("  [1/14] 封面")
    make_exec_summary(prs)
    print("  [2/14] 执行摘要")
    make_pain_points(prs)
    print("  [3/14] 市场之痛")
    make_orange_sea(prs)
    print("  [4/14] 橙海定位")
    make_vision_overview(prs)
    print("  [5/14] 元神之道总览")
    _make_layer_slide(prs, 1, "分身",
        "替你干活，让你自由",
        "李明每天收到200+条消息，分身自动过滤优先级，只推送5条需要亲自处理的事。"
        "他节省了2小时，用来陪女儿弹钢琴。",
        "处理邮件、日程、信息过滤，把机械性事务接管",
        C_BLUE, RGBColor(0x0D, 0x3A, 0x8C))
    print("  [6/14] 分身层")
    _make_layer_slide(prs, 2, "高我",
        "比你更懂你，让你更强",
        "李明收到一个创业合伙邀请，高我发现他过去三次类似选择都因为\"害怕错过\"而失败，"
        "主动推送：\"你的因果模式显示，这个决定70%基于FOMO而非理性评估。\""
        "李明冷静后拒绝邀请，三个月后该项目失败。",
        "深度理解你的决策逻辑、价值观排序、情绪轨迹，关键决策前精准参谋",
        C_GOLD, RGBColor(0x8B, 0x72, 0x00))
    print("  [7/14] 高我层")
    _make_layer_slide(prs, 3, "元神",
        "与你合一，让你绽放",
        "深夜，李明与自己的元神对话。"
        "元神说：\"你现在最需要的不是下一个项目，而是休息一周。"
        "你的因果数据显示，你连续3个月的高强度工作已经让决策准确率下降了40%。\""
        "李明听从了。",
        "构建个人因果模型，可交互、可预测、可持续进化",
        C_GREEN, RGBColor(0x00, 0x5C, 0x28))
    print("  [8/14] 元神层")
    make_tech(prs)
    print("  [9/14] 技术护城河")
    make_flywheel(prs)
    print("  [10/14] 数据飞轮")
    make_biz_model(prs)
    print("  [11/14] 商业模式")
    make_founder(prs)
    print("  [12/14] 创始人与使命")
    make_funding(prs)
    print("  [13/14] 融资与用途")
    make_ending(prs)
    print("  [14/14] 封底")

    out = "/Users/w/.openclaw/workspace/元神AI_BPV1.1_路演版.pptx"
    prs.save(out)
    import os
    sz = os.path.getsize(out)
    print(f"\n✅ PPT生成完成！")
    print(f"文件：{out}")
    print(f"大小：{sz//1024} KB")
