#!/usr/bin/env python3
"""
元神AI融资BP PPT生成脚本
基于SSS级视频框架提炼核心内容
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色方案
PRIMARY_BLUE = RGBColor(0, 82, 147)      # 深蓝
ACCENT_BLUE = RGBColor(0, 150, 255)      # 亮蓝
DARK_GRAY = RGBColor(51, 51, 51)          # 深灰
LIGHT_GRAY = RGBColor(128, 128, 128)     # 浅灰
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景形状
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_BLUE
    bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, content_items, has_numbers=False):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏背景
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_BLUE
    title_bg.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 内容区域
    y_pos = 1.8
    for i, item in enumerate(content_items):
        prefix = f"{i+1}. " if has_numbers else "• "
        text = prefix + item if has_numbers else item
        
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos), Inches(11.5), Inches(0.8))
        tf = content_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = DARK_GRAY
        p.line_spacing = 1.5
        
        y_pos += 0.9
    
    return slide

def add_data_slide(title, data_points):
    """添加数据幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_BLUE
    title_bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 数据卡片
    col = 0
    row = 0
    for label, value in data_points:
        x = 0.8 + col * 4
        y = 1.8 + row * 2.5
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.5), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(240, 248, 255)
        card.line.color.rgb = ACCENT_BLUE
        
        # 数值
        val_box = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.3), Inches(3.1), Inches(0.9))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # 标签
        lbl_box = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+1.2), Inches(3.1), Inches(0.6))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        col += 1
        if col >= 3:
            col = 0
            row += 1
    
    return slide

# ===== 开始生成幻灯片 =====

# 1. 封面
add_title_slide(
    "元神AI",
    "个人数字主权的觉醒系统 | 融资商业计划书"
)

# 2. 为什么是现在？
add_content_slide(
    "为什么是现在？",
    [
        "2026年，全球AI智能体市场规模61-99亿美元",
        "中国市场135亿元，增速超70%",
        "从'聊天时代'到'办事时代'的关键转折",
        "巨头观望，12-18个月窗口期",
        "这不是小风口，这是一个大时代"
    ]
)

# 3. 痛点定义
add_content_slide(
    "问题：数字记忆的碎片化",
    [
        "朋友圈、微博、知乎、小红书...你的足迹散落各处",
        "没有人能完整看到你的人生——包括你自己",
        "记忆数据被平台割裂，用户无法掌控",
        "这不是工具问题，这是主权问题"
    ]
)

# 4. 解决方案
add_content_slide(
    "元神AI：个人数字主权的觉醒系统",
    [
        "整合散落在各平台的记忆，还给你完整的人生",
        "基因胶囊：越用越强的系统，数据越多越聪明",
        "Agent Graph：智能体网络，连接所有平台",
        "从'替代'到'赋能'，让技术为你服务"
    ]
)

# 5. 市场规模
add_data_slide(
    "市场规模",
    [
        ("目标用户", "6000万人"),
        ("潜在收入", "60亿元"),
        ("年增速", "30-50%"),
        ("对标Notion", "100亿美元估值"),
        ("对标Obsidian", "10亿美元估值"),
        ("中国市场", "135亿元")
    ]
)

# 6. 竞争格局
add_content_slide(
    "为什么巨头不做？",
    [
        "隐私敏感：用户数据太敏感，巨头不敢碰",
        "非高频需求：不像聊天、邮件那样日常",
        "长期运营：不能快速变现",
        "这正是我们的机会——蓝海市场"
    ]
)

# 7. 三道壁垒
add_content_slide(
    "三道壁垒：为什么我们能赢？",
    [
        "1. 灵魂壁垒：十年四域淬炼的AI能力，5-10年积累不可复制",
        "2. 系统壁垒：基因胶囊正反馈循环，用户数据越多系统越强",
        "3. 生态壁垒：Agent Graph网络效应，连接所有平台"
    ],
    has_numbers=True
)

# 8. 商业模式
add_content_slide(
    "商业模式：按结果付费",
    [
        "66%企业倾向按业务成果计费",
        "56-65%用户愿为结果付费",
        "2028年70%软件供应商将重构商业模式",
        "免费基础功能 + 高级功能订阅 + 企业定制"
    ]
)

# 9. 财务预测
add_data_slide(
    "财务预测",
    [
        ("3个月", "1000用户"),
        ("6个月", "盈亏平衡"),
        ("12个月", "月收80万"),
        ("24个月", "年收1.25亿"),
        ("付费率", "30-50%"),
        ("用户规模", "50万")
    ]
)

# 10. 团队介绍
add_content_slide(
    "团队：十年四域淬炼的梦之队",
    [
        "创始人：十年AI科学+产业应用经验",
        "核心团队：技术+商业+国际视野",
        "顾问团队：行业专家+投资人",
        "元模型DNA：人与AI共处的深刻理解"
    ]
)

# 11. 融资需求
add_data_slide(
    "融资需求：76.5万启动",
    [
        ("团队", "45万"),
        ("基础设施", "8.75万"),
        ("运营", "9.5万"),
        ("预留", "12.65万"),
        ("总融资额", "76.5万"),
        ("盈亏平衡", "6个月")
    ]
)

# 12. 路线图
add_content_slide(
    "90天路线图",
    [
        "Phase 1（30天）：MVP开发完成",
        "Phase 2（15天）：种子用户内测",
        "Phase 3（15天）：功能扩展",
        "Phase 4（30天）：融资启动",
        "每个阶段都有明确的产出和风险控制"
    ],
    has_numbers=True
)

# 13. 结尾
add_title_slide(
    "投资一个时代",
    "元神AI | 让众生物质丰盈、精神富足、灵魂觉醒"
)

# 保存文件
output_path = "/Users/w/.openclaw/workspace/元神AI_融资BP.pptx"
prs.save(output_path)
print(f"✅ PPT已生成: {output_path}")
