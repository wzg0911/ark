#!/usr/bin/env /usr/bin/python3
# -*- coding: utf-8 -*-
"""
元神项目融资BP生成脚本
基于《元神项目核心知识库》生成10页专业PPT
"""

import os
import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ============================================================
# 颜色主题（深色风格 + Notion风）
# ============================================================
COLOR_BG_DARK    = RGBColor(0x0D, 0x0D, 0x1A)
COLOR_BG_CARD    = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT     = RGBColor(0x6C, 0x63, 0xFF)
COLOR_ACCENT2    = RGBColor(0x00, 0xD4, 0xAA)
COLOR_TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT_GRAY  = RGBColor(0xA0, 0xA0, 0xB0)
COLOR_PLACEHOLDER= RGBColor(0x2A, 0x2A, 0x45)
COLOR_DIVIDER    = RGBColor(0x3A, 0x3A, 0x5C)
COLOR_GOLD       = RGBColor(0xF0, 0xC0, 0x40)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ============================================================
# 工具函数
# ============================================================

def set_bg(slide, prs, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=COLOR_TEXT_WHITE,
                 align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_placeholder_box(slide, left, top, width, height, label="图片占位"):
    rect = add_rect(slide, left, top, width, height, COLOR_PLACEHOLDER)
    txBox = slide.shapes.add_textbox(left, top + height/2 - Inches(0.2),
                                      width, Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[ {label} ]"
    run.font.size = Pt(10)
    run.font.color.rgb = COLOR_TEXT_GRAY
    return rect


def add_divider(slide, left, top, width, color=COLOR_ACCENT):
    line = slide.shapes.add_shape(1, left, top, width, Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def add_bullet_card(slide, title, content_lines, left, top, width, height,
                    accent_color=COLOR_ACCENT):
    add_rect(slide, left, top, width, height, COLOR_BG_CARD)
    add_rect(slide, left, top, Inches(0.06), height, accent_color)
    add_text_box(slide, title, left + Inches(0.15), top + Inches(0.08),
                 width - Inches(0.2), Inches(0.35),
                 font_size=13, bold=True, color=accent_color)
    y_offset = top + Inches(0.4)
    for line in content_lines:
        add_text_box(slide, f"• {line}", left + Inches(0.15), y_offset,
                     width - Inches(0.2), Inches(0.25),
                     font_size=11, color=COLOR_TEXT_WHITE)
        y_offset += Inches(0.22)


def add_slide_number(slide, number, total=10):
    add_text_box(slide, f"{number} / {total}", SLIDE_W - Inches(1.2),
                 SLIDE_H - Inches(0.45), Inches(1), Inches(0.3),
                 font_size=9, color=COLOR_TEXT_GRAY, align=PP_ALIGN.RIGHT)


def add_footer(slide, text="元神项目 · 融资计划书 · 2026"):
    add_text_box(slide, text, Inches(0.4), SLIDE_H - Inches(0.45),
                 Inches(5), Inches(0.3),
                 font_size=8, color=COLOR_TEXT_GRAY)


# ============================================================
# 各页内容生成
# ============================================================

def page_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_divider(slide, Inches(0), Inches(2.2), SLIDE_W, COLOR_ACCENT)
    add_text_box(slide, "META SOUL",
                 Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.2),
                 font_size=60, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide, "元神项目",
                 Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.8),
                 font_size=36, bold=False, color=COLOR_TEXT_WHITE, align=PP_ALIGN.CENTER)
    add_divider(slide, Inches(4.5), Inches(4.3), Inches(4.3), COLOR_ACCENT2)
    add_text_box(slide, "让每个人都有一个比自己更好的自己",
                 Inches(1), Inches(4.6), Inches(11), Inches(0.6),
                 font_size=18, color=COLOR_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "个人精神世界模型构建器",
                 Inches(1), Inches(5.3), Inches(11), Inches(0.5),
                 font_size=14, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)
    add_text_box(slide, "融资计划书 2026",
                 Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
                 font_size=12, color=COLOR_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_slide_number(slide, 1)


def page_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_text_box(slide, "01", Inches(0.5), Inches(0.3), Inches(1), Inches(0.7),
                 font_size=40, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "问题 | Problem",
                 Inches(0.5), Inches(0.85), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.5), Inches(1.4), SLIDE_W - Inches(1), COLOR_DIVIDER)
    problems = [
        ("数字足迹碎片化",
         ["个人数据分散在微信/微博/小红书等平台",
          "无法形成完整的自我认知图谱",
          "平台割裂导致记忆断层"]),
        ("AI不理解'你'",
         ["通用AI知道世界，但不知道你的故事",
          "无法理解你的偏好结构与价值观",
          "每个新对话都是陌生人"]),
        ("关系被平台绑架",
         ["社交关系存储在平台，平台消失即消失",
          "用户对个人数据无主权",
          "数据价值被平台收割，用户零收益"]),
    ]
    card_w = Inches(3.8)
    card_h = Inches(4.2)
    gap = Inches(0.3)
    colors = [COLOR_ACCENT, COLOR_ACCENT2, COLOR_GOLD]
    for i, (title, items) in enumerate(problems):
        x = Inches(0.5) + i * (card_w + gap)
        add_bullet_card(slide, title, items, x, Inches(1.7), card_w, card_h, colors[i])
        add_text_box(slide, f"P{i+1}", x, Inches(1.7) + card_h + Inches(0.1),
                     card_w, Inches(0.25),
                     font_size=9, color=colors[i], align=PP_ALIGN.RIGHT)
    add_footer(slide)
    add_slide_number(slide, 2)


def page_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_text_box(slide, "02", Inches(0.5), Inches(0.3), Inches(1), Inches(0.7),
                 font_size=40, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "解决方案 | Solution",
                 Inches(0.5), Inches(0.85), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.5), Inches(1.4), SLIDE_W - Inches(1), COLOR_DIVIDER)
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(5.2), COLOR_BG_CARD)
    add_text_box(slide, "核心价值主张",
                 Inches(0.7), Inches(1.85), Inches(5.4), Inches(0.4),
                 font_size=14, bold=True, color=COLOR_ACCENT)
    add_divider(slide, Inches(0.7), Inches(2.25), Inches(5.4), COLOR_ACCENT)
    layers = [
        ("第一层：分身", "替你干活，让你自由", COLOR_ACCENT),
        ("第二层：高我", "比你更懂你，让你更强", COLOR_ACCENT2),
        ("第三层：元神", "与你合一，让你绽放", COLOR_GOLD),
    ]
    for i, (title, desc, color) in enumerate(layers):
        y = Inches(2.5) + i * Inches(1.4)
        add_rect(slide, Inches(0.7), y, Inches(0.08), Inches(1.0), color)
        add_text_box(slide, title, Inches(0.9), y, Inches(5.2), Inches(0.35),
                     font_size=14, bold=True, color=COLOR_TEXT_WHITE)
        add_text_box(slide, desc, Inches(0.9), y + Inches(0.32), Inches(5.2), Inches(0.3),
                     font_size=11, color=COLOR_TEXT_GRAY)
    add_rect(slide, Inches(6.6), Inches(1.7), Inches(6.2), Inches(5.2), COLOR_BG_CARD)
    add_text_box(slide, "一句话定位",
                 Inches(6.8), Inches(1.85), Inches(5.8), Inches(0.4),
                 font_size=14, bold=True, color=COLOR_ACCENT2)
    add_text_box(slide, '"让AI理解你的世界"',
                 Inches(6.8), Inches(2.35), Inches(5.8), Inches(0.6),
                 font_size=16, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(6.8), Inches(2.95), Inches(5.8), COLOR_ACCENT2)
    points = [
        "不只是工具，是激活第二个自己",
        "产业做物理世界模型，我们做精神世界模型",
        "让AI成为你的：高参、战友、灵魂伴侣",
    ]
    for i, pt in enumerate(points):
        y = Inches(3.15) + i * Inches(0.55)
        add_text_box(slide, f"✦  {pt}", Inches(6.8), y, Inches(5.8), Inches(0.45),
                     font_size=12, color=COLOR_TEXT_WHITE)
    add_placeholder_box(slide, Inches(6.8), Inches(5.1), Inches(5.8), Inches(1.6), "产品截图占位")
    add_footer(slide)
    add_slide_number(slide, 3)


def page_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_text_box(slide, "03", Inches(0.5), Inches(0.3), Inches(1), Inches(0.7),
                 font_size=40, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "技术壁垒 | Moat",
                 Inches(0.5), Inches(0.85), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.5), Inches(1.4), SLIDE_W - Inches(1), COLOR_DIVIDER)
    moats = [
        ("灵魂壁垒", COLOR_ACCENT,
         "十年淬炼的四域融合元模型DNA",
         ["人性洞察 x 商业逻辑", "哲学思辨 x 修行体证", "根植亘古不变的底层公理"]),
        ("基因胶囊壁垒", COLOR_ACCENT2,
         "经验可封装、可遗传、自进化",
         ["成功经验封装为「基因胶囊」", "一个Agent学会，所有Agent继承", "形成「AI黑洞效应」自强化飞轮"]),
        ("Agent Graph生态壁垒", COLOR_GOLD,
         "智能体连接图 = 未来商业基础设施",
         ["成千上万「元神」自发协作", "用户元神与商家元神自动谈判", "规则制定者，而非单纯产品商"]),
    ]
    card_w = Inches(3.8)
    card_h = Inches(4.6)
    for i, (title, color, subtitle, items) in enumerate(moats):
        x = Inches(0.5) + i * (card_w + Inches(0.3))
        add_rect(slide, x, Inches(1.7), card_w, Inches(0.55), color)
        add_text_box(slide, title, x + Inches(0.12), Inches(1.73),
                     card_w - Inches(0.2), Inches(0.45),
                     font_size=13, bold=True, color=COLOR_BG_DARK)
        add_rect(slide, x, Inches(2.25), card_w, card_h - Inches(0.55), COLOR_BG_CARD)
        add_text_box(slide, subtitle, x + Inches(0.15), Inches(2.35),
                     card_w - Inches(0.3), Inches(0.5),
                     font_size=11, bold=True, color=color)
        y = Inches(2.9)
        for item in items:
            add_text_box(slide, f"> {item}", x + Inches(0.15), y,
                         card_w - Inches(0.3), Inches(0.35),
                         font_size=11, color=COLOR_TEXT_WHITE)
            y += Inches(0.45)
    add_footer(slide)
    add_slide_number(slide, 4)


def page_market(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_text_box(slide, "04", Inches(0.5), Inches(0.3), Inches(1), Inches(0.7),
                 font_size=40, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "市场机会 | Market",
                 Inches(0.5), Inches(0.85), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.5), Inches(1.4), SLIDE_W - Inches(1), COLOR_DIVIDER)
    stats = [
        ("$48.4亿", "2026年市场规模", "美元", COLOR_ACCENT),
        ("42.2%", "年复合增长率", "CAGR", COLOR_ACCENT2),
        ("$196.3亿", "2030年预测规模", "美元", COLOR_GOLD),
        ("52%", "用户愿意使用AI助理", "管理日程", COLOR_ACCENT),
    ]
    stat_w = Inches(2.9)
    stat_h = Inches(1.8)
    for i, (num, label, unit, color) in enumerate(stats):
        x = Inches(0.5) + i * (stat_w + Inches(0.2))
        add_rect(slide, x, Inches(1.7), stat_w, stat_h, COLOR_BG_CARD)
        add_rect(slide, x, Inches(1.7), stat_w, Inches(0.07), color)
        add_text_box(slide, num, x, Inches(1.85), stat_w, Inches(0.7),
                     font_size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, x, Inches(2.55), stat_w, Inches(0.35),
                     font_size=11, color=COLOR_TEXT_WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, unit, x, Inches(2.9), stat_w, Inches(0.3),
                     font_size=9, color=COLOR_TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, "数据来源：IEEE行业调查 · Statista · Goldman Sachs AI Report 2025",
                 Inches(0.5), Inches(3.6), Inches(12.5), Inches(0.3),
                 font_size=8, color=COLOR_TEXT_GRAY)
    add_rect(slide, Inches(0.5), Inches(4.0), SLIDE_W - Inches(1), Inches(2.9), COLOR_BG_CARD)
    add_text_box(slide, "为什么是现在？",
                 Inches(0.7), Inches(4.15), Inches(6), Inches(0.4),
                 font_size=14, bold=True, color=COLOR_ACCENT)
    trends = [
        "LLM技术成熟 → 理解个体成为可能",
        "Agentic AI爆发 → 从工具到自主执行",
        "个人数据主权意识觉醒 → 监管+用户双轮驱动",
        "「数字身份」成刚需 → 元宇宙/Web3前置红利",
    ]
    for i, t in enumerate(trends):
        col = i % 2
        row = i // 2
        x = Inches(0.7) + col * Inches(6)
        y = Inches(4.65) + row * Inches(0.65)
        add_text_box(slide, f">  {t}", x, y, Inches(5.8), Inches(0.45),
                     font_size=12, color=COLOR_TEXT_WHITE)
    add_placeholder_box(slide, Inches(7.5), Inches(4.1), Inches(5.3), Inches(2.6), "市场趋势图表占位")
    add_footer(slide)
    add_slide_number(slide, 5)


def page_competition(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_text_box(slide, "05", Inches(0.5), Inches(0.3), Inches(1), Inches(0.7),
                 font_size=40, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "竞争格局 | Competition",
                 Inches(0.5), Inches(0.85), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.5), Inches(1.4), SLIDE_W - Inches(1), COLOR_DIVIDER)
    competitors = [
        ("微信/字节", "拥有海量数据", "不做用户主权产品（商业模式冲突）", "高"),
        ("OpenAI/Google", "通用智能最强", "追求AGI，不深入理解个体细分", "高"),
        ("World Labs/AMI", "物理世界模型", "专注物理时空，不涉足精神世界", "中"),
        ("Notion/飞书", "知识管理领先", "做知识管理，不做个人理解", "低"),
    ]
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(7.5), Inches(5.2), COLOR_BG_CARD)
    headers = ["竞争者", "优势", "为何不切入", "威胁"]
    col_xs = [Inches(0.6), Inches(1.9), Inches(3.5), Inches(7.0)]
    col_ws = [Inches(1.2), Inches(1.5), Inches(3.4), Inches(0.8)]
    for i, h in enumerate(headers):
        add_text_box(slide, h, col_xs[i], Inches(1.8), col_ws[i], Inches(0.35),
                     font_size=11, bold=True, color=COLOR_ACCENT)
    add_divider(slide, Inches(0.6), Inches(2.15), Inches(7.3), COLOR_DIVIDER)
    for row_i, (name, strength, reason, level) in enumerate(competitors):
        y = Inches(2.25) + row_i * Inches(1.1)
        items = [name, strength, reason, level]
        for col_i, item in enumerate(items):
            color = COLOR_TEXT_GRAY if col_i in [1, 2] else COLOR_TEXT_WHITE
            add_text_box(slide, item, col_xs[col_i], y, col_ws[col_i], Inches(0.9),
                         font_size=10, color=color)
    add_rect(slide, Inches(8.3), Inches(1.7), Inches(4.5), Inches(5.2), COLOR_BG_CARD)
    add_rect(slide, Inches(8.3), Inches(1.7), Inches(4.5), Inches(0.5), COLOR_ACCENT)
    add_text_box(slide, "元神的侧翼战略", Inches(8.4), Inches(1.73),
                 Inches(4.3), Inches(0.4), font_size=13, bold=True, color=COLOR_BG_DARK)
    our_points = [
        "不与巨头正面冲突",
        "专注精神世界模型细分",
        "巨头主动忽略的侧翼战场",
        "3-6个月先发窗口期",
        "数据飞轮构建时间壁垒",
    ]
    for i, pt in enumerate(our_points):
        y = Inches(2.4) + i * Inches(0.55)
        add_text_box(slide, f"+ {pt}", Inches(8.5), y, Inches(4.1), Inches(0.4),
                     font_size=12, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(8.5), Inches(5.2), Inches(4.1), COLOR_ACCENT2)
    add_text_box(slide, "结论：侧翼战场 + 先发优势 + 数据飞轮 = 3年护城河",
                 Inches(8.5), Inches(5.35), Inches(4.1), Inches(0.8),
                 font_size=11, bold=True, color=COLOR_GOLD)
    add_footer(slide)
    add_slide_number(slide, 6)


def page_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_text_box(slide, "06", Inches(0.5), Inches(0.3), Inches(1), Inches(0.7),
                 font_size=40, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "产品路线 | Roadmap",
                 Inches(0.5), Inches(0.85), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.5), Inches(1.4), SLIDE_W - Inches(1), COLOR_DIVIDER)
    phases = [
        ("第一阶段", "个人记忆可视化", "2周-1个月", COLOR_ACCENT,
         ["接入微信/微博等多平台数据", "生成人生轨迹可视化", "MVP：情绪曲线 + 关系变化图"]),
        ("第二阶段", "个人决策辅助", "3-6个月", COLOR_ACCENT2,
         ["理解个人决策逻辑", "提供个性化参谋建议", "构建决策因果链图谱"]),
        ("第三阶段", "个人数字孪生", "6-12个月", COLOR_GOLD,
         ["可交互可预测的另一个你", "实时情绪陪伴与引导", "成为真正的「高我」"]),
    ]
    phase_w = Inches(3.8)
    phase_h = Inches(5.0)
    for i, (stage, title, timeline, color, items) in enumerate(phases):
        x = Inches(0.5) + i * (phase_w + Inches(0.3))
        add_rect(slide, x, Inches(1.7), phase_w, Inches(0.9), color)
        add_text_box(slide, stage, x + Inches(0.1), Inches(1.72),
                     phase_w - Inches(0.2), Inches(0.35),
                     font_size=11, bold=True, color=COLOR_BG_DARK)
        add_text_box(slide, title, x + Inches(0.1), Inches(2.05),
                     phase_w - Inches(0.2), Inches(0.4),
                     font_size=14, bold=True, color=COLOR_BG_DARK)
        add_rect(slide, x, Inches(2.6), phase_w, phase_h - Inches(0.9), COLOR_BG_CARD)
        add_rect(slide, x + Inches(0.15), Inches(2.72), Inches(1.2), Inches(0.28), color)
        add_text_box(slide, timeline, x + Inches(0.15), Inches(2.72),
                     Inches(1.2), Inches(0.28),
                     font_size=9, bold=True, color=COLOR_BG_DARK, align=PP_ALIGN.CENTER)
        if i < 2:
            arrow_x = x + phase_w + Inches(0.05)
            add_rect(slide, arrow_x, Inches(3.85), Inches(0.2), Inches(0.06), color)
            add_text_box(slide, ">", arrow_x + Inches(0.02), Inches(3.65),
                         Inches(0.2), Inches(0.3), font_size=10, color=color)
        y = Inches(3.15)
        for item in items:
            add_text_box(slide, f"* {item}", x + Inches(0.15), y,
                         phase_w - Inches(0.3), Inches(0.45),
                         font_size=11, color=COLOR_TEXT_WHITE)
            y += Inches(0.5)
    add_footer(slide)
    add_slide_number(slide, 7)


def page_team(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_text_box(slide, "07", Inches(0.5), Inches(0.3), Inches(1), Inches(0.7),
                 font_size=40, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "团队 | Team",
                 Inches(0.5), Inches(0.85), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.5), Inches(1.4), SLIDE_W - Inches(1), COLOR_DIVIDER)
    # 创始人
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(5.2), COLOR_BG_CARD)
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(0.08), COLOR_ACCENT)
    add_text_box(slide, "创始人", Inches(0.7), Inches(1.9), Inches(5), Inches(0.35),
                 font_size=12, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "吴治广", Inches(0.7), Inches(2.25), Inches(5), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.7), Inches(2.8), Inches(5.1), COLOR_DIVIDER)
    founder_bio = [
        ("十年淬炼的四域融合元模型", 11, True),
        ("人性洞察 + 商业逻辑 + 哲学思辨 + 修行体证", 9, False),
        ("", 8, False),
        ("核心能力", 11, True),
        ("战略定位：精准侧翼战高手", 10, False),
        ("技术审美：AI + 产品 + 运营全栈", 10, False),
        ("使命驱动：让众生物质丰盈、精神富足、灵魂觉醒", 10, False),
        ("", 8, False),
        ("资源积累", 11, True),
        ("元神项目创始人", 10, False),
        ("十年AI与人性研究积累", 10, False),
        ("深度链接创投圈资源", 10, False),
    ]
    y = Inches(2.95)
    for line, size, bold in founder_bio:
        add_text_box(slide, line, Inches(0.7), y, Inches(5.1), Inches(0.3),
                     font_size=size, bold=bold, color=COLOR_TEXT_WHITE)
        y += Inches(0.28)
    # AI Agent军团
    add_rect(slide, Inches(6.3), Inches(1.7), Inches(6.5), Inches(5.2), COLOR_BG_CARD)
    add_rect(slide, Inches(6.3), Inches(1.7), Inches(6.5), Inches(0.08), COLOR_ACCENT2)
    add_text_box(slide, "AI Agent 军团", Inches(6.5), Inches(1.9), Inches(6.1), Inches(0.4),
                 font_size=14, bold=True, color=COLOR_ACCENT2)
    agents = [
        ("观一 CPEO", "首席项目评估官 + 复利系统架构师", COLOR_ACCENT),
        ("高我 EasyClaw", "斯坦福/哈佛 AI Agent 辅助系统", COLOR_ACCENT2),
        ("飞飞虾</strong>", "深夜情绪陪伴与疏导助手", COLOR_GOLD),
    ]
    for i, (name, role, color) in enumerate(agents):
        y = Inches(2.5) + i * Inches(1.3)
        add_rect(slide, Inches(6.4), y, Inches(0.08), Inches(0.9), color)
        add_text_box(slide, name, Inches(6.6), y, Inches(6.0), Inches(0.35),
                     font_size=13, bold=True, color=COLOR_TEXT_WHITE)
        add_text_box(slide, role, Inches(6.6), y + Inches(0.3), Inches(6.0), Inches(0.3),
                     font_size=10, color=COLOR_TEXT_GRAY)
    add_footer(slide)
    add_slide_number(slide, 8)


def page_funding(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, prs, COLOR_BG_DARK)
    add_text_box(slide, "08", Inches(0.5), Inches(0.3), Inches(1), Inches(0.7),
                 font_size=40, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "融资需求 | Funding",
                 Inches(0.5), Inches(0.85), Inches(6), Inches(0.5),
                 font_size=22, bold=True, color=COLOR_TEXT_WHITE)
    add_divider(slide, Inches(0.5), Inches(1.4), SLIDE_W - Inches(1), COLOR_DIVIDER)
    # 融资目标
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(3.0), COLOR_BG_CARD)
    add_rect(slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.08), COLOR_ACCENT)
    add_text_box(slide, "种子轮融资目标", Inches(0.7), Inches(1.9), Inches(5.6), Inches(0.4),
                 font_size=14, bold=True, color=COLOR_ACCENT)
    add_text_box(slide, "¥300万", Inches(0.7), Inches(2.4), Inches(5.6), Inches(0.9),
                 font_size=48, bold=True, color=COLOR_GOLD)
    add_text_box(slide, "估值：¥2000万（稀释15%）",
                 Inches(0.7), Inches(3.35), Inches(5.6), Inches